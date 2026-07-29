import argparse
import csv
import json
import os
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from textwrap import shorten

os.environ.setdefault("MPLCONFIGDIR", str(Path("results/matplotlib_cache")))

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from paper_plotting import (
    FALLBACK_COLORS as PAPER_FALLBACK_COLORS,
    MODEL_COLORS as PAPER_MODEL_COLORS,
    PAPER_FIGSIZE,
    finish_epoch_axis,
    get_model_style,
    mark_every,
    place_comparison_legend,
    save_figure_pair,
    setup_paper_plot_style,
)
from result_paths import build_report_artifact_dirs, resolve_run_artifact_paths


PRIORITY_METRICS = [
    "val_loss",
    "val_acc",
    "train_loss",
    "train_acc",
    "test_loss",
    "test_acc",
]
DEFAULT_CURVE_METRICS = ["val_loss", "val_acc", "train_loss", "train_acc"]
SUMMARY_METRIC_PRIORITY = [
    "test_macro_f1",
    "val_macro_f1",
    "test_macro_precision",
    "test_macro_recall",
    "val_macro_precision",
    "val_macro_recall",
]
PER_CLASS_PRIORITY = [
    "test_per_class_accuracy",
    "test_per_class_f1",
    "test_per_class_precision",
    "test_per_class_recall",
]
RUN_INFO_PRIORITY = [
    "model.architecture",
    "model.variant",
    "training.batch_size",
    "training.lr",
    "training.weight_decay",
    "training.epochs",
    "dataset.name",
    "dataset.image_size",
    "image_size",
    "batch_size",
    "lr",
    "weight_decay",
    "epochs",
    "weights",
    "embedding_dropout",
    "attention_dropout",
    "projection_dropout",
    "mlp_dropout",
    "early_stopping_metric",
    "early_stopping_patience",
    "val_ratio",
]
CONFIG_PRIORITY_PREFIXES = ["model.", "training.", "dataset.", "device", "command"]
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

COLOR_BG = (248, 250, 252)
COLOR_NAVY = (15, 23, 42)
COLOR_BLUE = (37, 99, 235)
COLOR_CYAN = (8, 145, 178)
COLOR_TEXT = (30, 41, 59)
COLOR_MUTED = (100, 116, 139)
COLOR_BORDER = (203, 213, 225)
COLOR_GOOD = (22, 163, 74)
COLOR_BAD = (220, 38, 38)
COLOR_CARD = (255, 255, 255)
COLOR_HEADER_FILL = (226, 232, 240)
FONT_FAMILY = "Aptos"
MODEL_DISPLAY_NAMES = {
    "vit_baseline": "ViT Baseline (No Pos)",
    "vit_learnable_position": "ViT Learnable Position",
    "vit_rope": "ViT RoPE",
    "vit_rope_2d": "ViT RoPE 2D",
    "vit_row_sinusoidal": "ViT Row-wise Sinusoidal",
    "vit_col_sinusoidal": "ViT Column-wise Sinusoidal",
    "vit_additive_sinusoidal": "ViT Additive Sinusoidal",
    "vit_additive_sinusoidal_shifted": "ViT Additive Sinusoidal Shifted",
    "vit_multiplicative_sinusoidal": "ViT Multiplicative Sinusoidal",
    "vit_multiplicative_sinusoidal_shifted": "ViT Multiplicative Sinusoidal Shifted",
    "resnet18_scratch": "ResNet18 Scratch",
    "resnet18_imagenet": "ResNet18 ImageNet",
    "vit": "ViT",
    "resnet18": "ResNet18",
}
MODEL_FAMILY_DISPLAY = {
    "vit": "ViT",
    "resnet18": "ResNet18",
    "cnn": "CNN",
}
MODEL_VARIANT_DISPLAY = {
    "baseline": "Baseline (No Pos)",
    "learnable_position": "Learnable Position",
    "rope": "RoPE",
    "rope_2d": "RoPE 2D",
    "row_sinusoidal": "Row-wise Sinusoidal",
    "col_sinusoidal": "Column-wise Sinusoidal",
    "additive_sinusoidal": "Additive Sinusoidal",
    "additive_sinusoidal_shifted": "Additive Sinusoidal Shifted",
    "multiplicative_sinusoidal": "Multiplicative Sinusoidal",
    "multiplicative_sinusoidal_shifted": "Multiplicative Sinusoidal Shifted",
    "scratch": "Scratch",
    "imagenet": "ImageNet",
    "pretrained": "Pretrained",
}
POSITION_ENCODING_DISPLAY = {
    "absolute": "Absolute",
    "rope": "RoPE",
    "rope_2d": "RoPE 2D",
    "row_sinusoidal": "Row-wise Sinusoidal",
    "col_sinusoidal": "Column-wise Sinusoidal",
    "additive_sinusoidal": "Additive Sinusoidal",
    "additive_sinusoidal_shifted": "Additive Sinusoidal Shifted",
    "multiplicative_sinusoidal": "Multiplicative Sinusoidal",
    "multiplicative_sinusoidal_shifted": "Multiplicative Sinusoidal Shifted",
    "none": "None",
}
INITIALIZATION_DISPLAY = {
    "scratch": "From scratch",
    "imagenet": "ImageNet pretrained",
    "none": "From scratch",
    "unknown": "Unknown",
}
DATASET_LABEL_NAMES = {
    "cadb_elements": [
        "horizontal",
        "vertical",
        "diagonal",
        "triangle",
        "symmetric",
        "pattern",
    ],
}
DATASET_DISPLAY_NAMES = {
    "cifar10": "CIFAR-10",
    "cadb_elements": "CADB Elements",
    "cadb_orientation": "CADB Orientation",
    "cadb_scene": "CADB Scene",
    "synthetic_orientation": "Synthetic Orientation",
    "synthetic_orientation_clean": "Synthetic Orientation Clean",
    "synthetic_orientation_hard": "Synthetic Orientation Hard",
    "synthetic_row_code": "Synthetic Row Code",
    "synthetic_col_code": "Synthetic Column Code",
}
COMPACT_CURVE_PRIORITY = ["val_macro_f1", "val_acc", "test_acc", "train_loss"]
COMPACT_PER_CLASS_PRIORITY = [
    "test_per_class_f1",
    "test_per_class_recall",
    "test_per_class_precision",
    "test_per_class_accuracy",
]
RUN_COLOR_MAP = dict(PAPER_MODEL_COLORS)
FALLBACK_RUN_COLORS = list(PAPER_FALLBACK_COLORS)


@dataclass
class ComparisonContext:
    report_name: str
    title: str
    scenario: str
    compact_mode: bool
    runs: list
    metrics: list[str]
    selected_metric_keys: list[str]
    macro_metric_keys: list[str]
    summary_rows: list[dict]
    varying_rows: list[tuple]
    all_varying_keys: list[str]
    metric_figure_paths: dict[str, Path]
    macro_figure_payload: dict | None
    per_class_figure_payloads: list[dict]
    headline_insights: list[str]
    conclusion_lines: list[str]


@dataclass
class RunArtifacts:
    run_name: str
    label: str
    history: list[dict]
    config: dict
    summary: dict
    available_metrics: list[str]
    flat_config: dict
    model_name: str
    model_family: str
    model_variant: str
    position_encoding: str
    initialization: str
    model_display_name: str
    device: str
    completed_epochs: int
    selected_epoch: int | None
    selected_metrics: dict[str, float]
    per_class_metrics: dict[str, list[float]]
    confusion_matrix_csv: Path | None
    confusion_matrix_figure: Path | None
    early_stopping: dict | None
    experiment_name: str | None
    results_experiment_dir: Path | None


def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate a comparison report and meeting-ready PPT for experiment runs."
    )
    input_group = parser.add_mutually_exclusive_group(required=True)
    input_group.add_argument(
        "--run",
        dest="runs",
        action="append",
        help="Run spec in the form run_name or run_name=Display Label. Repeatable.",
    )
    input_group.add_argument(
        "--summary-report",
        type=str,
        default=None,
        help="Existing seed-summary report folder name. If experiment-name is set, it is resolved inside that experiment first.",
    )
    input_group.add_argument(
        "--summary-manifest",
        type=Path,
        default=None,
        help="Path to a seed-summary summary_manifest.json file.",
    )
    parser.add_argument(
        "--per-class-report",
        type=str,
        default=None,
        help="Optional per-class comparison report folder name to append as extra slides.",
    )
    parser.add_argument(
        "--per-class-manifest",
        type=Path,
        default=None,
        help="Optional path to a per-class comparison report_manifest.json file.",
    )
    parser.add_argument("--results-dir", type=Path, default=Path("results"))
    parser.add_argument("--experiment-name", type=str, default=None)
    parser.add_argument("--report-name", type=str, default=None)
    parser.add_argument("--title", type=str, default=None)
    parser.add_argument(
        "--metrics",
        nargs="+",
        default=None,
        help="Metrics to compare. Defaults to the intersection of numeric metrics across runs.",
    )
    parser.add_argument(
        "--max-config-rows",
        type=int,
        default=18,
        help="Maximum number of varying config rows to show in the PPT run-info/config sections.",
    )
    parser.add_argument(
        "--skip-ppt",
        action="store_true",
        help="Skip PowerPoint generation and only write plots and summary files.",
    )
    parser.add_argument(
        "--compact",
        action="store_true",
        help="Export a shorter weekly-meeting PPT with fewer slides and only the highest-signal sections.",
    )
    return parser.parse_args()


def parse_run_spec(spec: str):
    if "=" in spec:
        run_name, label = spec.split("=", 1)
        return run_name.strip(), label.strip()
    return spec.strip(), spec.strip()


def load_json(path: Path):
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def load_csv_dict_rows(path: Path):
    with path.open("r", encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle))


def load_history(path: Path):
    rows = []
    with path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            parsed = {}
            for key, value in row.items():
                if value is None or value == "":
                    parsed[key] = value
                    continue
                try:
                    parsed[key] = float(value)
                except ValueError:
                    parsed[key] = value
            rows.append(parsed)
    return rows


def flatten_dict(data, prefix=""):
    flat = {}
    for key, value in data.items():
        full_key = f"{prefix}.{key}" if prefix else key
        if isinstance(value, dict):
            flat.update(flatten_dict(value, full_key))
        else:
            flat[full_key] = value
    return flat


def parse_summary_value(value):
    if value is None or value == "":
        return value
    try:
        numeric = float(value)
    except ValueError:
        return value
    if numeric.is_integer():
        return int(numeric)
    return numeric


def extract_command_model_name(command: str):
    tokens = str(command).split()
    if "--model" in tokens:
        index = tokens.index("--model")
        if index + 1 < len(tokens):
            return tokens[index + 1].strip()
    return None


def format_family_display(family: str):
    return MODEL_FAMILY_DISPLAY.get(family, family.replace("_", " ").title() if family else "Unknown")


def format_variant_display(variant: str):
    return MODEL_VARIANT_DISPLAY.get(variant, variant.replace("_", " ").title() if variant else "Unknown")


def format_position_encoding_display(position_encoding: str):
    return POSITION_ENCODING_DISPLAY.get(
        position_encoding,
        position_encoding.replace("_", " ").title() if position_encoding else "Unknown",
    )


def format_initialization_display(initialization: str):
    return INITIALIZATION_DISPLAY.get(
        initialization,
        initialization.replace("_", " ").title() if initialization else "Unknown",
    )


def format_dataset_display_name(dataset_name: str | None):
    if not dataset_name:
        return None
    return DATASET_DISPLAY_NAMES.get(dataset_name, dataset_name.replace("_", " ").title())


def build_dataset_plot_title(dataset_name: str | None, base_title: str):
    dataset_display_name = format_dataset_display_name(dataset_name)
    if not dataset_display_name:
        return base_title
    return f"{dataset_display_name} | {base_title}"


def build_model_display_name(model_name: str, family: str, variant: str, initialization: str):
    if model_name in MODEL_DISPLAY_NAMES:
        return MODEL_DISPLAY_NAMES[model_name]
    if family == "vit":
        if variant == "baseline":
            return "ViT Baseline (No Pos)"
        if variant == "learnable_position":
            return "ViT Learnable Position"
        if variant == "rope":
            return "ViT RoPE"
        return f"ViT {format_variant_display(variant)}"
    if family == "resnet18":
        if initialization == "imagenet" or variant == "imagenet":
            return "ResNet18 ImageNet"
        if initialization == "scratch" or variant == "scratch":
            return "ResNet18 Scratch"
        return f"ResNet18 {format_variant_display(variant)}"
    if family:
        return f"{format_family_display(family)} {format_variant_display(variant)}".strip()
    return model_name or "Unknown Model"


def infer_model_metadata(config: dict, summary: dict, run_name: str):
    selected = summary.get("selected_model", {}) if isinstance(summary.get("selected_model"), dict) else {}
    summary_config = summary.get("config", {}) if isinstance(summary.get("config"), dict) else {}
    model_cfg = config.get("model", {}) if isinstance(config.get("model"), dict) else {}

    command_model_name = extract_command_model_name(str(config.get("command", "")))
    raw_model_name = (
        selected.get("model_name")
        or summary_config.get("model_name")
        or summary_config.get("model")
        or command_model_name
    )
    family = (
        selected.get("model_family")
        or summary_config.get("model_family")
        or model_cfg.get("family")
        or model_cfg.get("architecture")
    )
    variant = (
        selected.get("model_variant")
        or summary_config.get("model_variant")
        or model_cfg.get("variant")
    )
    position_encoding = (
        selected.get("position_encoding")
        or summary_config.get("position_encoding")
        or model_cfg.get("position_encoding")
    )

    if not raw_model_name and family and variant:
        raw_model_name = f"{family}_{variant}"

    if not family and raw_model_name:
        if raw_model_name.startswith("vit"):
            family = "vit"
        elif raw_model_name.startswith("resnet18"):
            family = "resnet18"
    if not family:
        command = str(config.get("command", ""))
        if "train_cnn_cifar10.py" in command:
            family = "resnet18"
        elif "train_cifar10.py" in command or "train_cifar10_experiment.py" in command:
            family = "vit"
        elif config.get("weights") is not None:
            family = "resnet18"
        elif any(key in config for key in ["embedding_dropout", "attention_dropout", "mlp_dropout"]):
            family = "vit"

    if not variant and raw_model_name:
        if raw_model_name in {
            "vit_baseline",
            "vit_learnable_position",
            "vit_rope",
            "resnet18_scratch",
            "resnet18_imagenet",
        }:
            variant = raw_model_name.split("_", 1)[1]
        elif raw_model_name == "resnet18":
            variant = "baseline"
        elif raw_model_name == "vit":
            variant = "baseline"

    weights_value = (
        model_cfg.get("weights")
        or summary_config.get("weights")
        or config.get("weights")
    )
    if family == "resnet18":
        if str(weights_value).lower() == "imagenet" or variant in {"imagenet", "pretrained"}:
            initialization = "imagenet"
        elif str(weights_value).lower() == "none" or variant in {"scratch", "none"}:
            initialization = "scratch"
        else:
            initialization = "unknown"
    else:
        initialization = "scratch" if family == "vit" else "unknown"

    if not position_encoding:
        if variant == "rope":
            position_encoding = "rope"
        elif family == "vit":
            position_encoding = "absolute"
        else:
            position_encoding = "none"

    if not raw_model_name:
        if family and variant:
            raw_model_name = f"{family}_{variant}"
        elif family:
            raw_model_name = family
        else:
            raw_model_name = run_name

    display_name = build_model_display_name(raw_model_name, family or "unknown", variant or "unknown", initialization)
    return {
        "model_name": raw_model_name,
        "model_family": family or "unknown",
        "model_variant": variant or "unknown",
        "position_encoding": position_encoding or "unknown",
        "initialization": initialization or "unknown",
        "model_display_name": display_name,
    }


def resolve_optional_path(project_root: Path, value):
    if not value:
        return None
    path = Path(str(value))
    candidates = []
    if path.is_absolute():
        candidates.append(path)
    else:
        candidates.append(project_root / path)
        candidates.append((project_root / path.as_posix()).resolve())
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    return None


def extract_selected_metrics(summary: dict):
    selected = summary.get("selected_model")
    metrics = {}
    if not isinstance(selected, dict):
        return metrics
    for key, value in selected.items():
        if isinstance(value, (int, float)) and key != "epoch":
            metrics[key] = float(value)
    return metrics


def extract_per_class_metrics(summary: dict):
    selected = summary.get("selected_model")
    metrics = {}
    if not isinstance(selected, dict):
        return metrics
    for key, value in selected.items():
        if key.startswith("test_per_class_") and isinstance(value, list) and value:
            if all(isinstance(item, (int, float)) for item in value):
                metrics[key] = [float(item) for item in value]
    return metrics


def load_run_artifacts(
    results_dir: Path,
    project_root: Path,
    run_name: str,
    label: str,
    experiment_name: str | None = None,
):
    artifact_paths = resolve_run_artifact_paths(results_dir, run_name, experiment_name=experiment_name)
    history_path = artifact_paths["metrics_path"]
    config_path = artifact_paths["config_path"]
    summary_path = artifact_paths["summary_path"]

    required_paths = [history_path, config_path, summary_path]
    missing = [name for path, name in zip(required_paths, ["metrics.csv", "config.json", "summary.json"]) if path is None]
    if missing:
        raise FileNotFoundError(
            f"Missing artifacts for run '{run_name}': {', '.join(missing)}"
        )

    history = load_history(history_path)
    if not history:
        raise ValueError(f"No metric rows found in {history_path}")

    available_metrics = [
        key
        for key, value in history[0].items()
        if key != "epoch" and isinstance(value, (int, float))
    ]
    config = load_json(config_path)
    summary = load_json(summary_path)
    flat_config = flatten_dict(config)
    selected = summary.get("selected_model", {}) if isinstance(summary.get("selected_model"), dict) else {}
    model_metadata = infer_model_metadata(config, summary, run_name)
    effective_label = model_metadata["model_display_name"] if label == run_name else label

    return RunArtifacts(
        run_name=run_name,
        label=effective_label,
        history=history,
        config=config,
        summary=summary,
        available_metrics=available_metrics,
        flat_config=flat_config,
        model_name=model_metadata["model_name"],
        model_family=model_metadata["model_family"],
        model_variant=model_metadata["model_variant"],
        position_encoding=model_metadata["position_encoding"],
        initialization=model_metadata["initialization"],
        model_display_name=model_metadata["model_display_name"],
        device=str(config.get("device", "unknown")),
        completed_epochs=int(history[-1]["epoch"]),
        selected_epoch=int(selected["epoch"]) if isinstance(selected.get("epoch"), (int, float)) else None,
        selected_metrics=extract_selected_metrics(summary),
        per_class_metrics=extract_per_class_metrics(summary),
        confusion_matrix_csv=resolve_optional_path(
            project_root, selected.get("test_confusion_matrix_csv")
        ),
        confusion_matrix_figure=resolve_optional_path(
            project_root, selected.get("test_confusion_matrix_figure")
        ),
        early_stopping=summary.get("early_stopping") if isinstance(summary.get("early_stopping"), dict) else None,
        experiment_name=artifact_paths.get("experiment_name"),
        results_experiment_dir=artifact_paths.get("results_experiment_dir"),
    )


def order_metrics(metrics):
    seen = set()
    ordered = []
    for metric in PRIORITY_METRICS:
        if metric in metrics and metric not in seen:
            ordered.append(metric)
            seen.add(metric)
    for metric in metrics:
        if metric not in seen:
            ordered.append(metric)
            seen.add(metric)
    return ordered


def determine_metrics(runs, explicit_metrics=None):
    available_sets = [set(run.available_metrics) for run in runs]
    shared_metrics = set.intersection(*available_sets)
    if explicit_metrics:
        missing = [metric for metric in explicit_metrics if metric not in shared_metrics]
        if missing:
            raise ValueError(
                "These metrics are not present in every run: " + ", ".join(missing)
            )
        return order_metrics(explicit_metrics)
    default_metrics = [metric for metric in DEFAULT_CURVE_METRICS if metric in shared_metrics]
    if default_metrics:
        return default_metrics
    return order_metrics([metric for metric in runs[0].available_metrics if metric in shared_metrics])


def determine_selected_metric_keys(runs, history_metrics):
    selected_sets = [set(run.selected_metrics.keys()) for run in runs if run.selected_metrics]
    if len(selected_sets) != len(runs):
        return []

    shared = set.intersection(*selected_sets)
    shared = {
        key
        for key in shared
        if key not in history_metrics
        and not key.startswith("test_per_class_")
        and "confusion_matrix" not in key
    }

    def priority(key):
        if key in SUMMARY_METRIC_PRIORITY:
            return SUMMARY_METRIC_PRIORITY.index(key)
        return len(SUMMARY_METRIC_PRIORITY)

    return sorted(shared, key=lambda key: (priority(key), key))


def determine_macro_metric_keys(selected_metric_keys, metrics):
    available_metrics = set(selected_metric_keys) | {
        metric_name for metric_name in metrics if "macro_" in metric_name
    }
    preferred = [
        "test_macro_f1",
        "test_macro_precision",
        "test_macro_recall",
        "val_macro_f1",
        "val_macro_precision",
        "val_macro_recall",
    ]
    return [metric for metric in preferred if metric in available_metrics]


def detect_comparison_scenario(runs):
    families = {run.model_family for run in runs}
    variants = {run.model_variant for run in runs}

    if families == {"vit"} and variants == {"baseline", "row_sinusoidal", "col_sinusoidal"}:
        return "axis_bias_trio"
    if families == {"vit"} and variants == {"row_sinusoidal", "col_sinusoidal"}:
        return "row_vs_col_sinusoidal"
    if families == {"vit"} and {"baseline", "rope"}.issubset(variants):
        return "baseline_vs_rope"
    if families == {"resnet18"} and "scratch" in variants and ("imagenet" in variants or "pretrained" in variants):
        return "scratch_vs_pretrained"
    if "vit" in families and "resnet18" in families:
        return "vit_vs_cnn"
    return "generic_multi_run"


def build_default_title(runs, scenario):
    if scenario == "axis_bias_trio":
        return "Directional Comparison: ViT Baseline (No Pos) vs Row-wise vs Column-wise"
    if scenario == "row_vs_col_sinusoidal":
        return "Directional Comparison: ViT Row-wise vs ViT Column-wise"
    if scenario == "baseline_vs_rope":
        return "Weekly Comparison: ViT Baseline (No Pos) vs ViT RoPE"
    if scenario == "scratch_vs_pretrained":
        return "Weekly Comparison: ResNet18 Scratch vs ResNet18 ImageNet"
    if scenario == "vit_vs_cnn":
        return "Weekly Comparison: ViT vs CNN"
    joined = " vs ".join(run.label for run in runs[:3])
    return f"Experiment Comparison: {joined}" if joined else "Experiment Comparison Report"


def format_comparison_scenario(scenario: str):
    mapping = {
        "axis_bias_trio": "Baseline (No Pos) vs Row-wise vs Column-wise",
        "row_vs_col_sinusoidal": "Row-wise vs Column-wise Sinusoidal",
        "baseline_vs_rope": "ViT Baseline (No Pos) vs RoPE",
        "scratch_vs_pretrained": "ResNet18 Scratch vs Pretrained",
        "vit_vs_cnn": "ViT vs CNN",
        "generic_multi_run": "Multi-run comparison",
    }
    return mapping.get(scenario, scenario.replace("_", " ").title())


def ensure_unique_run_labels(runs):
    counts = {}
    for run in runs:
        counts[run.label] = counts.get(run.label, 0) + 1
    if all(count == 1 for count in counts.values()):
        return runs

    for run in runs:
        if counts[run.label] > 1:
            run.label = f"{run.label} ({run.run_name})"
    return runs


def is_percentage_metric(metric_name: str):
    tokens = ["acc", "accuracy", "precision", "recall", "f1"]
    lowered = metric_name.lower()
    return any(token in lowered for token in tokens)


def is_lower_better(metric_name: str):
    lowered = metric_name.lower()
    lower_keywords = ["loss", "error", "wer", "cer", "perplexity"]
    return any(keyword in lowered for keyword in lower_keywords)


def scale_metric_value(metric_name: str, value):
    if is_percentage_metric(metric_name):
        return value * 100.0
    return value


def metric_display_name(metric_name: str):
    display_names = {
        "train_loss": "Training Loss",
        "val_loss": "Validation Loss",
        "test_loss": "Test Loss",
        "train_acc": "Training Accuracy",
        "val_acc": "Validation Accuracy",
        "test_acc": "Test Accuracy",
    }
    if metric_name in display_names:
        return display_names[metric_name]
    tokens = metric_name.replace("_", " ").split()
    return " ".join(token.upper() if token in {"acc", "auc", "f1"} else token.capitalize() for token in tokens)


def format_metric_value(metric_name: str, value):
    scaled = scale_metric_value(metric_name, value)
    if is_percentage_metric(metric_name):
        return f"{scaled:.2f}%"
    return f"{scaled:.4f}"


def format_delta(metric_name: str, delta):
    if is_percentage_metric(metric_name):
        return f"{delta * 100.0:+.2f} pp"
    return f"{delta:+.4f}"


def compute_metric_summary(history, metric_name: str):
    rows = [row for row in history if metric_name in row]
    values = [row[metric_name] for row in rows]
    if is_lower_better(metric_name):
        best_index = min(range(len(values)), key=values.__getitem__)
    else:
        best_index = max(range(len(values)), key=values.__getitem__)

    final_row = rows[-1]
    best_row = rows[best_index]
    return {
        "final_value": final_row[metric_name],
        "final_epoch": int(final_row["epoch"]),
        "best_value": best_row[metric_name],
        "best_epoch": int(best_row["epoch"]),
    }


def build_summary_rows(runs, metrics, selected_metric_keys):
    rows = []
    for run in runs:
        row = {
            "run_name": run.run_name,
            "label": run.label,
            "model_name": run.model_name,
            "model_display_name": run.model_display_name,
            "model_family": run.model_family,
            "model_variant": run.model_variant,
            "position_encoding": run.position_encoding,
            "initialization": run.initialization,
            "device": run.device,
            "completed_epochs": run.completed_epochs,
            "selected_epoch": run.selected_epoch,
        }
        for metric in metrics:
            metric_summary = compute_metric_summary(run.history, metric)
            row[f"final_{metric}"] = metric_summary["final_value"]
            row[f"best_{metric}"] = metric_summary["best_value"]
            row[f"best_epoch_{metric}"] = metric_summary["best_epoch"]
        for metric_name in selected_metric_keys:
            if metric_name in run.selected_metrics:
                row[f"selected_{metric_name}"] = run.selected_metrics[metric_name]
        rows.append(row)
    return rows


def choose_varying_config_keys(runs, max_rows):
    all_keys = sorted(set().union(*(run.flat_config.keys() for run in runs)))
    varying = []
    for key in all_keys:
        values = [run.flat_config.get(key) for run in runs]
        normalized = [json.dumps(value, sort_keys=True, ensure_ascii=False) for value in values]
        if len(set(normalized)) > 1:
            varying.append(key)

    def priority(key):
        for index, prefix in enumerate(CONFIG_PRIORITY_PREFIXES):
            if key.startswith(prefix):
                return index
        return len(CONFIG_PRIORITY_PREFIXES)

    varying.sort(key=lambda key: (priority(key), key))
    return varying[:max_rows], varying


def stringify_config_value(value):
    if value is None:
        return "None"
    if isinstance(value, float):
        return f"{value:.6g}"
    if isinstance(value, (list, tuple)):
        return ", ".join(str(item) for item in value)
    return str(value)


def infer_shared_experiment_name(runs):
    experiment_names = {run.experiment_name for run in runs if run.experiment_name}
    if len(experiment_names) == 1:
        return next(iter(experiment_names))
    return None


def ensure_report_dirs(results_dir: Path, report_name: str, experiment_name: str | None = None, runs=None):
    resolved_experiment_name = experiment_name or infer_shared_experiment_name(runs or [])
    report_paths = build_report_artifact_dirs(
        results_dir=results_dir,
        report_name=report_name,
        experiment_name=resolved_experiment_name,
    )
    report_dir = report_paths["report_dir"]
    figures_dir = report_paths["figures_dir"]
    figures_dir.mkdir(parents=True, exist_ok=True)
    return report_dir, figures_dir


def write_csv(path: Path, fieldnames, rows):
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


def build_metric_insight_rows(summary_rows, metrics, reference_label):
    rows = []
    for metric in metrics:
        if is_lower_better(metric):
            best_row = min(summary_rows, key=lambda row: row[f"final_{metric}"])
        else:
            best_row = max(summary_rows, key=lambda row: row[f"final_{metric}"])
        reference_row = next((row for row in summary_rows if row["label"] == reference_label), summary_rows[0])
        delta = best_row[f"final_{metric}"] - reference_row[f"final_{metric}"]
        rows.append(
            {
                "metric": metric,
                "best_run": best_row["label"],
                "best_value": best_row[f"final_{metric}"],
                "reference_delta": delta,
            }
        )
    return rows


def get_shared_dataset_name(runs):
    dataset_names = set()
    for run in runs:
        dataset_value = run.config.get("dataset", run.flat_config.get("dataset"))
        if isinstance(dataset_value, dict):
            dataset_value = dataset_value.get("name")
        if dataset_value:
            dataset_names.add(str(dataset_value))
    if len(dataset_names) == 1:
        return next(iter(dataset_names))
    return None


def build_shared_dataset_card_content(runs):
    dataset_name = get_shared_dataset_name(runs)
    if not dataset_name or not runs:
        return None

    reference_run = runs[0]
    dataset_config = reference_run.config.get("dataset", {})
    cadb_config = reference_run.config.get("cadb_dataset", {})
    dataset_sizes = []
    for key in ("train_size", "val_size", "test_size"):
        value = dataset_config.get(key) if isinstance(dataset_config, dict) else None
        if value is not None:
            dataset_sizes.append(f"{key.replace('_size', '')} {value}")
    size_line = f"Sizes: {', '.join(dataset_sizes)}" if dataset_sizes else None

    if dataset_name == "cadb_elements":
        val_ratio = reference_run.flat_config.get("val_ratio", reference_run.config.get("training", {}).get("val_ratio"))
        lines = [
            "Task: CADB composition elements multi-label classification",
            "Files: images/ + composition_elements.json + split.json",
            "Labels: horizontal / vertical / diagonal",
            "triangle / symmetric / pattern",
            f"Split: official train/test + val from train (ratio={val_ratio})",
        ]
        if size_line:
            lines.append(size_line)
        return "CADB Dataset Structure", lines[:6]

    if dataset_name == "cadb_scene":
        lines = [
            "Task: CADB scene classification",
            "Files: images/ + scene_categories.json + split.json",
            "Labels: 10 scene categories",
            "Split: official train/test + val from train",
        ]
        if size_line:
            lines.append(size_line)
        return "CADB Dataset Structure", lines[:6]

    if dataset_name == "cadb_orientation":
        label_mode = cadb_config.get("label_mode", reference_run.flat_config.get("cadb_label_mode", "exclusive"))
        lines = [
            "Task: CADB horizontal-vs-vertical pilot subset",
            "Files: images/ + composition_elements.json + split.json",
            f"Label mode: {label_mode}",
            "Split: official train/test when available + val from train",
        ]
        if size_line:
            lines.append(size_line)
        return "CADB Dataset Structure", lines[:6]

    return None


def build_headline_insights(runs, summary_rows, metrics, scenario):
    if not metrics:
        return []

    reference_row = summary_rows[0]
    primary_metric = metrics[0]
    dataset_name = get_shared_dataset_name(runs)
    if is_lower_better(primary_metric):
        best_row = min(summary_rows, key=lambda row: row[f"final_{primary_metric}"])
    else:
        best_row = max(summary_rows, key=lambda row: row[f"final_{primary_metric}"])

    insights = []

    if scenario == "axis_bias_trio":
        insights.append(
            "This deck keeps the ViT backbone fixed and compares three positional choices: learned absolute embedding, row-tied sinusoidal indexing, and column-tied sinusoidal indexing.",
        )
        insights.append(
            f"Primary comparison metric: {metric_display_name(primary_metric)}. "
            f"Best final result is {best_row['label']} with "
            f"{format_metric_value(primary_metric, best_row[f'final_{primary_metric}'])}."
        )
        if dataset_name == "synthetic_orientation_clean":
            insights.append(
                "On the clean v2 synthetic split, all three models approach saturation, so the main value is as a sanity-check that the axis-biased variants do not break the task."
            )
            if len(summary_rows) >= 3:
                col_row = next((row for row in summary_rows if row["model_variant"] == "col_sinusoidal"), None)
                base_row = next((row for row in summary_rows if row["model_variant"] == "baseline"), None)
                if col_row and base_row:
                    delta = col_row[f"selected_test_macro_f1"] - base_row[f"selected_test_macro_f1"]
                    insights.append(
                        f"Column-wise reaches the cleanest selected-checkpoint result here, at {format_delta('test_macro_f1', delta)} versus the baseline on test macro F1."
                    )
        else:
            insights.append(
                "Use this three-way setup to separate the effect of adding an axis-tied bias from the effect of switching between row-tied and column-tied indexing."
            )
    else:
        insights.append(
            f"Primary comparison metric: {metric_display_name(primary_metric)}. "
            f"Best final result is {best_row['label']} with "
            f"{format_metric_value(primary_metric, best_row[f'final_{primary_metric}'])}."
        )

    if scenario == "row_vs_col_sinusoidal":
        insights.insert(
            0,
            "This deck isolates directional positional bias: the ViT backbone is unchanged, and only the fixed sinusoidal axis used for patch positions is swapped.",
        )
        insights.append(
            "Structurally, the row-wise variant gives patches in the same row the same positional code, while the column-wise variant ties patches in the same column."
        )
        if len(summary_rows) >= 2:
            row_run = next((row for row in summary_rows if row["model_variant"] == "row_sinusoidal"), None)
            col_run = next((row for row in summary_rows if row["model_variant"] == "col_sinusoidal"), None)
            if row_run and col_run:
                insights.append(
                    "Row-wise finishes slightly higher on final test accuracy, while Column-wise reaches the stronger best validation accuracy and selected test macro F1."
                )
    elif scenario == "baseline_vs_rope":
        insights.insert(
            0,
            "This deck is framed as a ViT baseline vs RoPE comparison, so position encoding is the first story to highlight.",
        )
    elif scenario == "scratch_vs_pretrained":
        insights.insert(
            0,
            "This deck is framed as a scratch vs pretrained comparison, so initialization should be called out before raw metric differences.",
        )
    elif scenario == "vit_vs_cnn":
        insights.insert(
            0,
            "This deck is framed as a ViT vs CNN comparison, so architecture family differences should be separated from training hyperparameter differences.",
        )

    if best_row["label"] != reference_row["label"]:
        delta = best_row[f"final_{primary_metric}"] - reference_row[f"final_{primary_metric}"]
        insights.append(
            (
                f"Against reference run {reference_row['label']}, {best_row['label']} changes "
                f"{metric_display_name(primary_metric)} by {format_delta(primary_metric, delta)}."
            )
        )

    if scenario != "axis_bias_trio" and "train_acc" in metrics and "test_acc" in metrics:
        gap_rows = []
        for row in summary_rows:
            gap = row["final_train_acc"] - row["final_test_acc"]
            gap_rows.append((gap, row["label"]))
        largest_gap, run_label = max(gap_rows)
        insights.append(
            (
                f"Largest train-to-test accuracy gap appears in {run_label} "
                f"({largest_gap * 100.0:.2f} pp), worth calling out during generalization discussion."
            )
        )

    stopped = []
    for run in runs:
        status = run.early_stopping or {}
        if status.get("enabled"):
            flag = "stopped early" if status.get("stopped_early") else "ran to planned epochs"
            stopped.append(f"{run.label}: {flag}")
    if stopped:
        insights.append("Early stopping summary: " + "; ".join(stopped[:3]))

    return insights[:4]


def build_conclusion_lines(runs, summary_rows, metrics, per_class_metric_keys, scenario):
    conclusions = []
    dataset_name = get_shared_dataset_name(runs)
    if metrics:
        primary_metric = metrics[0]
        if is_lower_better(primary_metric):
            winner = min(summary_rows, key=lambda row: row[f"final_{primary_metric}"])
        else:
            winner = max(summary_rows, key=lambda row: row[f"final_{primary_metric}"])
        conclusions.append(
            (
                f"Current headline result: {winner['label']} leads on "
                f"{metric_display_name(primary_metric)} with "
                f"{format_metric_value(primary_metric, winner[f'final_{primary_metric}'])}."
            )
        )

    if len(summary_rows) >= 2 and metrics and scenario not in {"row_vs_col_sinusoidal", "axis_bias_trio"}:
        reference = summary_rows[0]
        challenger = summary_rows[1]
        metric = metrics[0]
        delta = challenger[f"final_{metric}"] - reference[f"final_{metric}"]
        conclusions.append(
            (
                f"When presenting baseline vs comparison, use {reference['label']} as reference and report "
                f"{challenger['label']} at {format_delta(metric, delta)} on {metric_display_name(metric)}."
            )
        )

    if scenario == "axis_bias_trio":
        if dataset_name == "synthetic_orientation_clean":
            base_run = next((row for row in summary_rows if row["model_variant"] == "baseline"), None)
            col_run = next((row for row in summary_rows if row["model_variant"] == "col_sinusoidal"), None)
            if base_run and col_run:
                delta = col_run["selected_test_macro_f1"] - base_run["selected_test_macro_f1"]
                conclusions.append(
                    f"On clean v2, all three models essentially solve the task, but ViT Column-wise achieves the cleanest selected-checkpoint result at {format_delta('test_macro_f1', delta)} versus the baseline on test macro F1."
                )
        conclusions.append(
            "The structural sweep here is controlled: the attention blocks stay the same, and only the positional encoding changes from learned absolute to row-tied or column-tied fixed sinusoidal indexing."
        )
        if dataset_name == "synthetic_orientation_clean":
            conclusions.append(
                "Because the clean v2 split is already near-saturated, use this deck as a sanity-check slide: it shows the directional bias variants behave differently, but the task is too easy to be the main evidence slide."
            )
        else:
            conclusions.append(
                "Use this trio to discuss whether any gain comes from introducing an explicit axis bias at all, before narrowing to a row-vs-col comparison."
            )
    elif scenario == "row_vs_col_sinusoidal":
        row_run = next((row for row in summary_rows if row["model_variant"] == "row_sinusoidal"), None)
        col_run = next((row for row in summary_rows if row["model_variant"] == "col_sinusoidal"), None)
        if row_run and col_run:
            delta = col_run["final_test_acc"] - row_run["final_test_acc"]
            conclusions.append(
                f"On final test accuracy, ViT Column-wise trails ViT Row-wise by {format_delta('test_acc', delta)}, but Column-wise reaches the better best validation accuracy and selected test macro F1."
            )
        conclusions.append(
            "The structural change here is narrow and controlled: learned 2D-style absolute embeddings are replaced by fixed sinusoidal embeddings defined on only one spatial axis."
        )
        conclusions.append(
            "Use this comparison to discuss whether CADB orientation cues benefit more from row-tied or column-tied positional indexing, rather than from a larger model change."
        )
    elif scenario == "baseline_vs_rope":
        conclusions.append(
            "For baseline vs RoPE slides, keep the takeaway centered on positional encoding rather than on generic hyperparameter drift."
        )
    elif scenario == "scratch_vs_pretrained":
        conclusions.append(
            "For scratch vs pretrained slides, explicitly state whether the observed gain looks like an initialization effect or a persistent optimization gap."
        )
    elif scenario == "vit_vs_cnn":
        conclusions.append(
            "For ViT vs CNN slides, treat architecture family as the main axis and keep variant-level details as secondary notes."
        )

    if per_class_metric_keys:
        conclusions.append(
            "Per-class pages are now included when selected-checkpoint classwise metrics exist, so class imbalance or hard classes can be discussed directly in meetings."
        )

    conclusions.append(
        "The deck is now organized for weekly reporting: setup, overview, key metrics, curves, error analysis, and a conclusion slide instead of raw log-style export."
    )
    return conclusions[:4]


def chunked(items, chunk_size):
    for index in range(0, len(items), chunk_size):
        yield items[index : index + chunk_size]


def estimate_table_font_size(column_count, row_count):
    if column_count <= 5 and row_count <= 6:
        return 14
    if column_count <= 6 and row_count <= 8:
        return 12
    if column_count <= 8 and row_count <= 10:
        return 11
    return 10


def save_summary_outputs(
    report_dir: Path,
    title: str,
    scenario: str,
    runs,
    metrics,
    selected_metric_keys,
    summary_rows,
    varying_keys,
    headline_insights,
    conclusion_lines,
):
    summary_csv_path = report_dir / "comparison_summary.csv"
    fieldnames = [
        "run_name",
        "label",
        "model_name",
        "model_display_name",
        "model_family",
        "model_variant",
        "position_encoding",
        "initialization",
        "device",
        "completed_epochs",
        "selected_epoch",
    ]
    for metric in metrics:
        fieldnames.extend(
            [f"final_{metric}", f"best_{metric}", f"best_epoch_{metric}"]
        )
    for metric in selected_metric_keys:
        fieldnames.append(f"selected_{metric}")
    write_csv(summary_csv_path, fieldnames, summary_rows)

    config_csv_path = report_dir / "config_comparison.csv"
    config_rows = []
    for key in varying_keys:
        row = {"config_key": key}
        for run in runs:
            row[run.label] = stringify_config_value(run.flat_config.get(key))
        config_rows.append(row)
    write_csv(
        config_csv_path,
        ["config_key"] + [run.label for run in runs],
        config_rows,
    )

    overview_path = report_dir / "overview.md"
    overview_lines = [f"# {title}", ""]
    overview_lines.append(f"Generated: {datetime.now().isoformat(timespec='seconds')}")
    overview_lines.append(f"Comparison scenario: `{scenario}`")
    overview_lines.append(f"Reference run: `{runs[0].label}`")
    overview_lines.append("")
    overview_lines.append("## Headline Takeaways")
    overview_lines.append("")
    for line in headline_insights:
        overview_lines.append(f"- {line}")
    overview_lines.append("")
    overview_lines.append("## Runs")
    overview_lines.append("")
    for run in runs:
        overview_lines.append(
            f"- `{run.label}` (`{run.run_name}`), model: `{run.model_display_name}`, "
            f"family: `{run.model_family}`, variant: `{run.model_variant}`, "
            f"position encoding: `{run.position_encoding}`, epochs completed: `{run.completed_epochs}`, device: `{run.device}`"
        )
    overview_lines.append("")
    overview_lines.append("## Key Metrics")
    overview_lines.append("")
    for row in summary_rows:
        overview_lines.append(f"### {row['label']}")
        for metric in metrics:
            overview_lines.append(
                f"- {metric_display_name(metric)}: final {format_metric_value(metric, row[f'final_{metric}'])}, "
                f"best {format_metric_value(metric, row[f'best_{metric}'])} at epoch {row[f'best_epoch_{metric}']}"
            )
        for metric in selected_metric_keys:
            overview_lines.append(
                f"- Selected {metric_display_name(metric)}: {format_metric_value(metric, row[f'selected_{metric}'])}"
            )
        overview_lines.append("")
    overview_lines.append("## Suggested Meeting Conclusion")
    overview_lines.append("")
    for line in conclusion_lines:
        overview_lines.append(f"- {line}")
    overview_path.write_text("\n".join(overview_lines), encoding="utf-8")

    presentation_summary_path = report_dir / "presentation_summary.json"
    presentation_summary = {
        "title": title,
        "scenario": scenario,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "reference_run": {"run_name": runs[0].run_name, "label": runs[0].label},
        "headline_insights": headline_insights,
        "conclusion_lines": conclusion_lines,
        "metrics": metrics,
        "selected_checkpoint_metrics": selected_metric_keys,
        "runs": summary_rows,
    }
    presentation_summary_path.write_text(
        json.dumps(presentation_summary, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    manifest_path = report_dir / "report_manifest.json"
    manifest = {
        "report_dir": str(report_dir),
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "scenario": scenario,
        "runs": [{"run_name": run.run_name, "label": run.label} for run in runs],
        "metrics": metrics,
        "selected_checkpoint_metrics": selected_metric_keys,
        "summary_csv": str(summary_csv_path),
        "config_csv": str(config_csv_path),
        "overview_md": str(overview_path),
        "presentation_summary_json": str(presentation_summary_path),
    }
    manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    return summary_csv_path, config_csv_path, overview_path, presentation_summary_path, manifest_path


def resolve_summary_manifest_path(
    results_dir: Path,
    summary_report: str | None,
    summary_manifest: Path | None,
    experiment_name: str | None = None,
):
    if summary_manifest:
        path = summary_manifest if summary_manifest.is_absolute() else Path.cwd() / summary_manifest
        if path.is_dir():
            path = path / "summary_manifest.json"
    elif summary_report:
        candidates = []
        if experiment_name:
            candidates.append(results_dir / experiment_name / "reports" / summary_report / "summary_manifest.json")
        candidates.append(results_dir / "reports" / summary_report / "summary_manifest.json")
        recursive_matches = sorted(results_dir.glob(f"**/{summary_report}/summary_manifest.json"))
        candidates.extend(recursive_matches)
        path = next((candidate for candidate in candidates if candidate.exists()), candidates[0])
    else:
        raise ValueError("Either --summary-report or --summary-manifest must be provided.")

    if not path.exists():
        raise FileNotFoundError(f"Seed-summary manifest not found: {path}")
    return path.resolve()


def resolve_report_manifest_path(
    results_dir: Path,
    report_name: str | None,
    manifest_path: Path | None,
    experiment_name: str | None = None,
):
    if manifest_path:
        path = manifest_path if manifest_path.is_absolute() else Path.cwd() / manifest_path
        if path.is_dir():
            path = path / "report_manifest.json"
    elif report_name:
        candidates = []
        if experiment_name:
            candidates.append(results_dir / experiment_name / "reports" / report_name / "report_manifest.json")
        candidates.append(results_dir / "reports" / report_name / "report_manifest.json")
        recursive_matches = sorted(results_dir.glob(f"**/{report_name}/report_manifest.json"))
        candidates.extend(recursive_matches)
        path = next((candidate for candidate in candidates if candidate.exists()), candidates[0])
    else:
        return None

    if not path.exists():
        raise FileNotFoundError(f"Report manifest not found: {path}")
    return path.resolve()


def resolve_summary_artifact_path(project_root: Path, value):
    if not value:
        return None
    path = Path(str(value))
    candidates = [path] if path.is_absolute() else [project_root / path]
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    raise FileNotFoundError(f"Referenced summary artifact not found: {value}")


def format_class_triplet(items):
    return ", ".join(f"{name} ({value})" for name, value in items)


def mean_column(rows, key):
    values = [float(row[key]) for row in rows if row.get(key) not in (None, "")]
    return sum(values) / len(values) if values else float("-inf")


def top_class_deltas(rows, delta_key, reverse=True, top_k=3):
    values = []
    for row in rows:
        if delta_key not in row:
            continue
        values.append((float(row[delta_key]), row["class_name"]))
    values.sort(reverse=reverse)
    return [(name, f"{value * 100.0:+.2f} pp") for value, name in values[:top_k]]


def lowest_class_values(rows, value_key, top_k=3):
    values = []
    for row in rows:
        if value_key not in row:
            continue
        values.append((float(row[value_key]), row["class_name"]))
    values.sort()
    return [(name, f"{value * 100.0:.2f}%") for value, name in values[:top_k]]


def detect_seed_summary_scenario(model_names):
    families = set()
    variants = set()
    for model_name in model_names:
        if model_name.startswith("vit"):
            families.add("vit")
        elif model_name.startswith("resnet18"):
            families.add("resnet18")
        if "_" in model_name:
            variants.add(model_name.split("_", 1)[1])

    if families == {"vit"} and {"baseline", "rope"}.issubset(variants):
        return "baseline_vs_rope"
    if families == {"resnet18"} and "scratch" in variants and ("imagenet" in variants or "pretrained" in variants):
        return "scratch_vs_pretrained"
    if "vit" in families and "resnet18" in families:
        return "vit_vs_cnn"
    return "generic_multi_run"


def build_seed_summary_conclusion_lines(summary_rows, reference_model, metrics):
    if not summary_rows or not metrics:
        return []

    metric_priority = ["test_acc", "macro_f1", "best_val_acc"]
    primary_metric = next((metric for metric in metric_priority if metric in metrics), metrics[0])
    conclusions = []
    best_row = max(summary_rows, key=lambda item: item[f"{primary_metric}_mean"])
    conclusions.append(
        f"Current aggregate winner: {best_row['model_label']} leads on {metric_display_name(primary_metric)} "
        f"with {format_metric_value(primary_metric, best_row[f'{primary_metric}_mean'])}."
    )

    reference_row = next((row for row in summary_rows if row["model"] == reference_model), None)
    if reference_row is not None:
        for row in summary_rows:
            if row["model"] == reference_model:
                continue
            delta = row[f"{primary_metric}_mean"] - reference_row[f"{primary_metric}_mean"]
            conclusions.append(
                f"Relative to {reference_row['model_label']}, {row['model_label']} changes "
                f"{metric_display_name(primary_metric)} by {format_delta(primary_metric, delta)}."
            )

    conclusions.append(
        "This deck is based on multi-seed aggregate statistics, so the meeting takeaway should emphasize average behavior and variance rather than a single lucky run."
    )
    return conclusions[:4]


def save_seed_summary_report_outputs(
    report_dir: Path,
    title: str,
    scenario: str,
    source_manifest_path: Path,
    source_manifest: dict,
    headline_insights,
    conclusion_lines,
    per_class_manifest_path: Path | None = None,
):
    presentation_summary_path = report_dir / "presentation_summary.json"
    presentation_summary = {
        "title": title,
        "scenario": scenario,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "source_manifest": str(source_manifest_path),
        "headline_insights": headline_insights,
        "conclusion_lines": conclusion_lines,
        "report_type": "seed_summary",
        "summary_report_name": source_manifest.get("report_name"),
        "metrics": source_manifest.get("metrics", []),
        "models": source_manifest.get("models", []),
        "seeds": source_manifest.get("seeds", []),
        "reference_model": source_manifest.get("reference_model"),
        "per_class_report_manifest": str(per_class_manifest_path) if per_class_manifest_path else None,
    }
    presentation_summary_path.write_text(
        json.dumps(presentation_summary, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    manifest_path = report_dir / "report_manifest.json"
    manifest = {
        "report_dir": str(report_dir),
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "report_type": "seed_summary",
        "scenario": scenario,
        "title": title,
        "source_summary_manifest": str(source_manifest_path),
        "summary_report_name": source_manifest.get("report_name"),
        "metrics": source_manifest.get("metrics", []),
        "models": source_manifest.get("models", []),
        "seeds": source_manifest.get("seeds", []),
        "reference_model": source_manifest.get("reference_model"),
        "per_class_report_manifest": str(per_class_manifest_path) if per_class_manifest_path else None,
        "presentation_summary_json": str(presentation_summary_path),
    }
    manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    return presentation_summary_path, manifest_path


def load_per_class_report_data(args, project_root: Path):
    manifest_path = resolve_report_manifest_path(
        args.results_dir,
        args.per_class_report,
        args.per_class_manifest,
        experiment_name=args.experiment_name,
    )
    if manifest_path is None:
        return None

    manifest = load_json(manifest_path)
    figures = {}
    for key, value in (manifest.get("figures") or {}).items():
        figures[key] = resolve_summary_artifact_path(project_root, value)

    accuracy_csv_path = resolve_summary_artifact_path(project_root, manifest.get("accuracy_csv"))
    f1_csv_path = resolve_summary_artifact_path(project_root, manifest.get("f1_csv"))
    accuracy_rows = [
        {key: parse_summary_value(value) for key, value in row.items()}
        for row in load_csv_dict_rows(accuracy_csv_path)
    ]
    f1_rows = [
        {key: parse_summary_value(value) for key, value in row.items()}
        for row in load_csv_dict_rows(f1_csv_path)
    ]

    runs = manifest.get("runs", [])
    reference_run = str(manifest.get("reference_run", runs[0]["run_name"] if runs else "reference"))
    run_labels = {item["run_name"]: item.get("label", item["run_name"]) for item in runs}

    value_run_names = [
        item["run_name"]
        for item in runs
        if accuracy_rows and f"{item['run_name']}_value" in accuracy_rows[0]
    ]
    if not value_run_names:
        return None

    best_accuracy_run = max(value_run_names, key=lambda run_name: mean_column(accuracy_rows, f"{run_name}_value"))
    best_f1_run = max(value_run_names, key=lambda run_name: mean_column(f1_rows, f"{run_name}_value"))
    focus_run = next((run_name for run_name in value_run_names if run_name != reference_run), best_accuracy_run)
    focus_accuracy_run = best_accuracy_run if best_accuracy_run != reference_run else focus_run
    focus_f1_run = best_f1_run if best_f1_run != reference_run else focus_run

    accuracy_delta_key = f"{focus_accuracy_run}_delta_vs_{reference_run}"
    f1_delta_key = f"{focus_f1_run}_delta_vs_{reference_run}"

    slides = []
    if "per_class_accuracy_grouped" in figures:
        slides.append(
            {
                "title": "Per-Class Analysis | Accuracy",
                "subtitle": "Classwise test accuracy grouped across compared models.",
                "path": figures["per_class_accuracy_grouped"],
                "summary": [
                    f"Best mean class accuracy in this comparison: {run_labels.get(best_accuracy_run, best_accuracy_run)}.",
                    f"Top gains vs {run_labels.get(reference_run, reference_run)}: {format_class_triplet(top_class_deltas(accuracy_rows, accuracy_delta_key, reverse=True))}.",
                    f"Weakest classes in {run_labels.get(focus_accuracy_run, focus_accuracy_run)}: {format_class_triplet(lowest_class_values(accuracy_rows, f'{focus_accuracy_run}_value'))}.",
                ],
            }
        )

    if "per_class_accuracy_delta_vs_reference" in figures:
        slides.append(
            {
                "title": "Per-Class Analysis | Accuracy Delta",
                "subtitle": f"Accuracy delta is measured against {run_labels.get(reference_run, reference_run)}.",
                "path": figures["per_class_accuracy_delta_vs_reference"],
                "summary": [
                    f"Largest positive deltas for {run_labels.get(focus_accuracy_run, focus_accuracy_run)}: {format_class_triplet(top_class_deltas(accuracy_rows, accuracy_delta_key, reverse=True))}.",
                    f"Largest negative deltas: {format_class_triplet(top_class_deltas(accuracy_rows, accuracy_delta_key, reverse=False))}.",
                    "Use this page to show that class-level gains are directional rather than uniform across all categories.",
                ],
            }
        )

    if "per_class_f1_grouped" in figures:
        slides.append(
            {
                "title": "Per-Class Analysis | F1",
                "subtitle": "Classwise F1 grouped across compared models.",
                "path": figures["per_class_f1_grouped"],
                "summary": [
                    f"Best mean class F1 in this comparison: {run_labels.get(best_f1_run, best_f1_run)}.",
                    f"Top F1 gains vs {run_labels.get(reference_run, reference_run)}: {format_class_triplet(top_class_deltas(f1_rows, f1_delta_key, reverse=True))}.",
                    f"Weakest F1 classes in {run_labels.get(focus_f1_run, focus_f1_run)}: {format_class_triplet(lowest_class_values(f1_rows, f'{focus_f1_run}_value'))}.",
                ],
            }
        )

    return {
        "manifest_path": manifest_path,
        "manifest": manifest,
        "slides": slides,
    }


def setup_plot_style():
    setup_paper_plot_style()


def get_run_color(run, index: int):
    if run.model_name in RUN_COLOR_MAP:
        return RUN_COLOR_MAP[run.model_name]
    return FALLBACK_RUN_COLORS[index % len(FALLBACK_RUN_COLORS)]


def plot_metric(figures_dir: Path, runs, metric_name: str):
    setup_plot_style()
    figure, axis = plt.subplots(figsize=PAPER_FIGSIZE)
    dataset_name = get_shared_dataset_name(runs)
    for index, run in enumerate(runs):
        style = get_model_style(run.model_name, index)
        epochs = [int(row["epoch"]) for row in run.history]
        values = [scale_metric_value(metric_name, row[metric_name]) for row in run.history]
        axis.plot(
            epochs,
            values,
            marker=style["marker"],
            markevery=mark_every(len(epochs)),
            linestyle=style["linestyle"],
            linewidth=2.0,
            markersize=3.8,
            label=run.label,
            color=style["color"],
        )

    finish_epoch_axis(
        axis,
        metric_name=metric_name,
        title=build_dataset_plot_title(dataset_name, metric_display_name(metric_name)),
        show_legend=False,
    )
    place_comparison_legend(axis, len(runs))

    figure_path = figures_dir / f"{metric_name}_comparison.png"
    save_figure_pair(figure, figure_path)
    plt.close(figure)
    return figure_path


def write_publication_selected_checkpoint_table(report_dir: Path, runs, summary_rows, selected_metric_keys):
    table_path = report_dir / "publication_selected_checkpoints.csv"
    preferred_metrics = [
        "test_acc",
        "test_loss",
        "test_macro_f1",
        "val_acc",
        "val_loss",
        "val_macro_f1",
    ]
    metric_keys = [metric for metric in preferred_metrics if metric in selected_metric_keys]
    fieldnames = ["label", "run_name", "selected_epoch"] + [f"selected_{metric}" for metric in metric_keys]
    rows = []
    for row in summary_rows:
        output_row = {
            "label": row["label"],
            "run_name": row["run_name"],
            "selected_epoch": row["selected_epoch"],
        }
        for metric in metric_keys:
            value = row.get(f"selected_{metric}")
            output_row[f"selected_{metric}"] = format_metric_value(metric, value) if value is not None else ""
        rows.append(output_row)
    write_csv(table_path, fieldnames, rows)
    return table_path


def plot_selected_test_accuracy_summary(figures_dir: Path, runs, summary_rows):
    if not all("selected_test_acc" in row for row in summary_rows):
        return None

    setup_plot_style()
    sorted_rows = sorted(summary_rows, key=lambda row: row["selected_test_acc"])
    labels = [row["label"] for row in sorted_rows]
    values = [row["selected_test_acc"] * 100.0 for row in sorted_rows]
    y_positions = np.arange(len(sorted_rows))

    height = max(3.2, 0.42 * len(sorted_rows) + 1.2)
    figure, axis = plt.subplots(figsize=(7.2, height))
    colors = [
        get_model_style(
            next((run.model_name for run in runs if run.label == row["label"]), row["model_name"]),
            index,
        )["color"]
        for index, row in enumerate(sorted_rows)
    ]
    axis.scatter(values, y_positions, s=38, color=colors, zorder=3)
    axis.hlines(y_positions, xmin=0.0, xmax=values, color="#d1d5db", linewidth=1.0, zorder=1)
    for value, y_position in zip(values, y_positions):
        axis.text(value + 0.35, y_position, f"{value:.2f}%", va="center", fontsize=8.5)

    axis.set_yticks(y_positions)
    axis.set_yticklabels(labels)
    axis.set_xlabel("Selected Test Accuracy (%)")
    lower_bound = max(0.0, min(values) - 1.0)
    upper_bound = min(100.0, max(values) + 1.0)
    axis.set_xlim(lower_bound, upper_bound)
    axis.set_title("Selected-Checkpoint Test Accuracy", pad=10)
    axis.grid(True, axis="x", linestyle="--", linewidth=0.7, alpha=0.35)
    axis.grid(False, axis="y")
    axis.spines["top"].set_visible(False)
    axis.spines["right"].set_visible(False)

    figure_path = figures_dir / "paper_selected_test_accuracy.png"
    save_figure_pair(figure, figure_path)
    plt.close(figure)
    return figure_path


def write_publication_captions(report_dir: Path, title: str, runs, metrics, selected_metric_keys, publication_paths):
    dataset_name = format_dataset_display_name(get_shared_dataset_name(runs)) or "the dataset"
    run_count = len(runs)
    selected_list = ", ".join(metric_display_name(metric) for metric in selected_metric_keys[:4])
    captions_path = report_dir / "figure_captions.md"

    lines = [
        f"# Figure Captions: {title}",
        "",
        "These captions are drafts for thesis figures or supervisor updates. Edit the final wording after the final multi-seed rerun.",
        "",
    ]
    figure_number = 1
    for metric in metrics:
        metric_note = " Accuracy-style metrics are reported as percentages." if is_percentage_metric(metric) else ""
        lines.extend(
            [
                f"## Figure {figure_number}. {metric_display_name(metric)} across compared models",
                "",
                (
                    f"{metric_display_name(metric)} on {dataset_name} for {run_count} compared models. "
                    "All curves use the same split, optimizer settings, and checkpoint-selection rule within this report."
                    f"{metric_note}"
                ),
                "",
            ]
        )
        figure_number += 1
    if publication_paths.get("selected_test_accuracy"):
        lines.extend(
            [
                f"## Figure {figure_number}. Selected-checkpoint test performance",
                "",
                (
                    f"Test accuracy of each model at the checkpoint selected by validation performance on {dataset_name}. "
                    "This figure summarizes final held-out performance and should be interpreted together with validation loss "
                    "to separate peak accuracy from overfitting behaviour."
                ),
                "",
            ]
        )
        figure_number += 1
    if publication_paths.get("selected_checkpoint_table"):
        lines.extend(
            [
                "## Table 1. Selected-checkpoint metrics",
                "",
                (
                    f"Selected-checkpoint metrics for all compared models. "
                    f"The table includes selected epoch and shared selected metrics"
                    f"{': ' + selected_list if selected_list else ''}."
                ),
                "",
            ]
        )
    captions_path.write_text("\n".join(lines), encoding="utf-8")
    return captions_path


def plot_macro_metrics(figures_dir: Path, runs, macro_metric_keys):
    if not macro_metric_keys:
        return None
    dataset_name = get_shared_dataset_name(runs)

    def resolve_macro_value(run, metric_name: str):
        if metric_name in run.selected_metrics:
            return run.selected_metrics[metric_name]
        if metric_name in run.available_metrics:
            return compute_metric_summary(run.history, metric_name)["final_value"]
        raise KeyError(metric_name)

    focus_metrics = [
        metric
        for metric in ["test_macro_f1", "test_macro_precision", "test_macro_recall"]
        if metric in macro_metric_keys
    ]
    if not focus_metrics:
        focus_metrics = macro_metric_keys[:3]

    setup_plot_style()
    figure, axis = plt.subplots(figsize=(8.0, 4.8))
    x_positions = np.arange(len(focus_metrics))
    width = min(0.75 / max(1, len(runs)), 0.22)

    for index, run in enumerate(runs):
        color = get_run_color(run, index)
        values = [scale_metric_value(metric_name, resolve_macro_value(run, metric_name)) for metric_name in focus_metrics]
        offset = (index - (len(runs) - 1) / 2) * width
        axis.bar(
            x_positions + offset,
            values,
            width=width,
            label=run.label,
            color=color,
            alpha=0.9,
        )

    axis.set_xticks(x_positions)
    axis.set_xticklabels([metric_display_name(metric_name) for metric_name in focus_metrics], fontsize=10)
    axis.set_ylabel("Percentage (%)")
    axis.set_ylim(0, max(100.0, axis.get_ylim()[1]))
    axis.set_title(build_dataset_plot_title(dataset_name, "Macro Metric Snapshot"), fontsize=14, pad=12)
    axis.legend(loc="upper center", ncol=min(3, len(runs)), frameon=True)
    figure.tight_layout()

    figure_path = figures_dir / "macro_metrics_comparison.png"
    figure.savefig(figure_path, dpi=180)
    plt.close(figure)

    summary_lines = []
    if "test_macro_f1" in focus_metrics:
        best_run = max(runs, key=lambda run: resolve_macro_value(run, "test_macro_f1"))
        best_value = resolve_macro_value(best_run, "test_macro_f1")
        summary_lines.append(
            f"Best Test Macro F1 snapshot: {best_run.label} ({format_metric_value('test_macro_f1', best_value)})."
        )
    if "test_macro_precision" in focus_metrics:
        best_precision_run = max(runs, key=lambda run: resolve_macro_value(run, "test_macro_precision"))
        summary_lines.append(
            f"Best selected Test Macro Precision: {best_precision_run.label} ({format_metric_value('test_macro_precision', resolve_macro_value(best_precision_run, 'test_macro_precision'))})."
        )
    if "test_macro_recall" in focus_metrics:
        best_recall_run = max(runs, key=lambda run: resolve_macro_value(run, "test_macro_recall"))
        summary_lines.append(
            f"Best selected Test Macro Recall: {best_recall_run.label} ({format_metric_value('test_macro_recall', resolve_macro_value(best_recall_run, 'test_macro_recall'))})."
        )

    return {
        "path": figure_path,
        "metrics": focus_metrics,
        "summary_lines": summary_lines[:3],
    }


def load_confusion_matrix_csv(path: Path):
    with path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle)
        rows = list(reader)
    if not rows or len(rows) < 2:
        raise ValueError(f"Confusion matrix CSV is empty: {path}")

    class_names = rows[0][1:]
    matrix = []
    true_labels = []
    for row in rows[1:]:
        true_labels.append(row[0])
        matrix.append([int(value) for value in row[1:]])
    return true_labels or class_names, matrix


def ensure_confusion_matrix_figure(figures_dir: Path, run: RunArtifacts):
    if run.confusion_matrix_csv and run.confusion_matrix_csv.exists():
        class_names, matrix = load_confusion_matrix_csv(run.confusion_matrix_csv)
    elif run.confusion_matrix_figure and run.confusion_matrix_figure.exists():
        return run.confusion_matrix_figure
    else:
        return None

    data = np.array(matrix, dtype=float)

    setup_plot_style()
    figure, axis = plt.subplots(figsize=(6.2, 5.4))
    heatmap = axis.imshow(data, cmap="Blues")
    axis.set_xticks(range(len(class_names)))
    axis.set_yticks(range(len(class_names)))
    axis.set_xticklabels(class_names, rotation=45, ha="right", fontsize=8)
    axis.set_yticklabels(class_names, fontsize=8)
    axis.set_xlabel("Predicted label")
    axis.set_ylabel("True label")
    dataset_name = run.config.get("dataset")
    if isinstance(dataset_name, dict):
        dataset_name = dataset_name.get("name")
    axis.set_title(
        build_dataset_plot_title(dataset_name, f"{run.model_display_name} Test Confusion Matrix"),
        fontsize=13,
        pad=10,
    )

    for row_index in range(data.shape[0]):
        for column_index in range(data.shape[1]):
            value = int(data[row_index, column_index])
            text_color = "white" if value > data.max() * 0.55 else "#0f172a"
            axis.text(column_index, row_index, str(value), ha="center", va="center", fontsize=7, color=text_color)

    figure.colorbar(heatmap, ax=axis, fraction=0.046, pad=0.04)
    figure.tight_layout()
    figure_path = figures_dir / f"{run.run_name}_confusion_matrix.png"
    figure.savefig(figure_path, dpi=180)
    plt.close(figure)
    return figure_path


def determine_available_per_class_metric_keys(runs):
    metric_sets = [set(run.per_class_metrics.keys()) for run in runs if run.per_class_metrics]
    if len(metric_sets) != len(runs):
        return []

    shared = set.intersection(*metric_sets)

    def priority(metric_name):
        if metric_name in PER_CLASS_PRIORITY:
            return PER_CLASS_PRIORITY.index(metric_name)
        return len(PER_CLASS_PRIORITY)

    return sorted(shared, key=lambda metric_name: (priority(metric_name), metric_name))


def infer_class_names(run: RunArtifacts, metric_name: str):
    if run.confusion_matrix_csv and run.confusion_matrix_csv.exists():
        class_names, _ = load_confusion_matrix_csv(run.confusion_matrix_csv)
        return class_names
    selected_model = run.summary.get("selected_model", {})
    if isinstance(selected_model, dict):
        label_names = selected_model.get("label_names")
        if isinstance(label_names, list) and label_names:
            return [str(label) for label in label_names]
    dataset_name = run.config.get("dataset")
    if isinstance(dataset_name, dict):
        dataset_name = dataset_name.get("name")
    if dataset_name in DATASET_LABEL_NAMES:
        return list(DATASET_LABEL_NAMES[dataset_name])
    value_count = len(run.per_class_metrics.get(metric_name, []))
    return [f"class_{index}" for index in range(value_count)]


def plot_per_class_metric(figures_dir: Path, runs, metric_name: str):
    dataset_name = get_shared_dataset_name(runs)
    class_names = infer_class_names(runs[0], metric_name)
    value_count = len(class_names)
    x_positions = np.arange(value_count)
    width = min(0.72 / max(1, len(runs)), 0.22)

    setup_plot_style()
    figure, axis = plt.subplots(figsize=(9.4, 4.8))

    for index, run in enumerate(runs):
        color = get_run_color(run, index)
        values = np.array(run.per_class_metrics[metric_name]) * 100.0
        offset = (index - (len(runs) - 1) / 2) * width
        axis.bar(
            x_positions + offset,
            values,
            width=width,
            label=run.label,
            color=color,
            alpha=0.9,
        )

    axis.set_xticks(x_positions)
    axis.set_xticklabels(class_names, rotation=35, ha="right", fontsize=8)
    axis.set_ylim(0, max(100.0, axis.get_ylim()[1]))
    axis.set_ylabel("Percentage (%)")
    axis.set_title(
        build_dataset_plot_title(dataset_name, f"{metric_display_name(metric_name)} by Class"),
        fontsize=14,
        pad=12,
    )
    axis.legend(loc="upper center", ncol=min(3, len(runs)), frameon=True)
    figure.tight_layout()

    figure_path = figures_dir / f"{metric_name}_per_class.png"
    figure.savefig(figure_path, dpi=180)
    plt.close(figure)
    return figure_path, class_names


def summarize_per_class_metric(runs, metric_name: str, class_names):
    mean_scores = []
    for run in runs:
        values = run.per_class_metrics[metric_name]
        mean_scores.append((sum(values) / len(values), run.label))
    mean_scores.sort(reverse=not is_lower_better(metric_name))

    per_class_spread = []
    for index, class_name in enumerate(class_names):
        class_values = [run.per_class_metrics[metric_name][index] for run in runs]
        spread = max(class_values) - min(class_values)
        per_class_spread.append((spread, class_name))
    per_class_spread.sort(reverse=True)

    lines = [
        f"Highest mean {metric_display_name(metric_name)}: {mean_scores[0][1]} ({mean_scores[0][0] * 100.0:.2f}%).",
        f"Largest cross-run spread appears on class '{per_class_spread[0][1]}' ({per_class_spread[0][0] * 100.0:.2f} pp).",
    ]
    return lines


def top_confusion_pairs(confusion_matrix_csv: Path):
    class_names, matrix = load_confusion_matrix_csv(confusion_matrix_csv)
    pairs = []
    for row_index, true_label in enumerate(class_names):
        for column_index, pred_label in enumerate(class_names):
            if row_index == column_index:
                continue
            count = matrix[row_index][column_index]
            if count > 0:
                pairs.append((count, true_label, pred_label))
    pairs.sort(reverse=True)
    return pairs[:3]


def build_varying_rows(runs, varying_keys, max_rows):
    rows = []
    for key in varying_keys[:max_rows]:
        values = [stringify_config_value(run.flat_config.get(key)) for run in runs]
        rows.append((key, values))
    return rows


def make_default_report_name(run_specs):
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    short_names = "_vs_".join(spec[0] for spec in run_specs[:2])
    short_names = short_names[:48]
    return f"comparison_{short_names}_{timestamp}"


def build_run_profile_lines(run: RunArtifacts, selected_varying_keys):
    lines = [
        f"Run ID: {run.run_name}",
        f"Model / device: {run.model_display_name} / {run.device}",
        f"Epochs completed: {run.completed_epochs}",
    ]

    lines.append(
        f"Family / variant: {format_family_display(run.model_family)} / {format_variant_display(run.model_variant)}"
    )
    lines.append(
        f"Position / init: {format_position_encoding_display(run.position_encoding)} / {format_initialization_display(run.initialization)}"
    )

    batch_size = run.flat_config.get("training.batch_size", run.flat_config.get("batch_size"))
    lr = run.flat_config.get("training.lr", run.flat_config.get("lr"))
    weight_decay = run.flat_config.get("training.weight_decay", run.flat_config.get("weight_decay"))
    if batch_size is not None or lr is not None or weight_decay is not None:
        lines.append(
            "Batch / LR / WD: "
            f"{stringify_config_value(batch_size)} / {stringify_config_value(lr)} / {stringify_config_value(weight_decay)}"
        )

    early_stopping = run.early_stopping or {}
    if early_stopping.get("enabled"):
        lines.append(
            "Early stopping: "
            f"{early_stopping.get('metric', 'unknown')} (patience={early_stopping.get('patience', 'n/a')}, "
            f"stopped_early={early_stopping.get('stopped_early', False)})"
        )

    return lines[:6]


def build_model_comparison_rows(runs):
    headers = [
        "Run",
        "Model",
        "Family",
        "Variant",
        "Position Encoding",
        "Initialization",
        "Batch / LR / WD",
    ]
    rows = []
    for run in runs:
        batch_size = run.flat_config.get("training.batch_size", run.flat_config.get("batch_size"))
        lr = run.flat_config.get("training.lr", run.flat_config.get("lr"))
        weight_decay = run.flat_config.get("training.weight_decay", run.flat_config.get("weight_decay"))
        rows.append(
            [
                run.label,
                run.model_display_name,
                format_family_display(run.model_family),
                format_variant_display(run.model_variant),
                format_position_encoding_display(run.position_encoding),
                format_initialization_display(run.initialization),
                f"{stringify_config_value(batch_size)} / {stringify_config_value(lr)} / {stringify_config_value(weight_decay)}",
            ]
        )
    return headers, rows


def build_macro_metric_rows(summary_rows, macro_metric_keys):
    headers = ["Macro Metric", "Direction"] + [row["label"] for row in summary_rows] + ["Best Run"]
    rows = []
    for metric_name in macro_metric_keys:
        row_key = f"selected_{metric_name}"
        if not all(row_key in row for row in summary_rows):
            row_key = f"final_{metric_name}"
        if not all(row_key in row for row in summary_rows):
            continue
        best_row = max(summary_rows, key=lambda row: row[row_key])
        rows.append(
            [
                f"{'Selected' if row_key.startswith('selected_') else 'Final'} {metric_display_name(metric_name)}",
                "Higher",
                *[format_metric_value(metric_name, row[row_key]) for row in summary_rows],
                best_row["label"],
            ]
        )
    return headers, rows


def build_results_overview_rows(summary_rows, metrics):
    chosen_metrics = metrics[: min(3, len(metrics))]
    headers = ["Run", "Model", "Variant", "Epochs", "Selected Epoch"]
    headers.extend(metric_display_name(metric) for metric in chosen_metrics)

    rows = []
    for row in summary_rows:
        rows.append(
            [
                row["label"],
                row["model_display_name"],
                format_variant_display(row["model_variant"]),
                str(row["completed_epochs"]),
                str(row["selected_epoch"]) if row["selected_epoch"] is not None else "n/a",
                *[format_metric_value(metric, row[f"final_{metric}"]) for metric in chosen_metrics],
            ]
        )
    return headers, rows


def build_main_metric_rows(summary_rows, metrics, selected_metric_keys):
    headers = ["Metric", "Direction"] + [row["label"] for row in summary_rows] + ["Best Run"]
    rows = []

    for metric in metrics:
        best_row = min(summary_rows, key=lambda row: row[f"final_{metric}"]) if is_lower_better(metric) else max(summary_rows, key=lambda row: row[f"final_{metric}"])
        rows.append(
            [
                f"Final {metric_display_name(metric)}",
                "Lower" if is_lower_better(metric) else "Higher",
                *[format_metric_value(metric, row[f"final_{metric}"]) for row in summary_rows],
                best_row["label"],
            ]
        )

    for metric in selected_metric_keys:
        best_row = min(summary_rows, key=lambda row: row[f"selected_{metric}"]) if is_lower_better(metric) else max(summary_rows, key=lambda row: row[f"selected_{metric}"])
        rows.append(
            [
                f"Selected {metric_display_name(metric)}",
                "Lower" if is_lower_better(metric) else "Higher",
                *[format_metric_value(metric, row[f"selected_{metric}"]) for row in summary_rows],
                best_row["label"],
            ]
        )

    return headers, rows


def build_compact_metric_rows(summary_rows):
    headers = ["Metric", "Direction"] + [row["label"] for row in summary_rows] + ["Best Run"]
    metric_specs = [
        ("Final Test ACC", "final_test_acc", "test_acc"),
        ("Final Val ACC", "final_val_acc", "val_acc"),
        ("Final Val Macro F1", "final_val_macro_f1", "val_macro_f1"),
        ("Final Test Macro F1", "final_test_macro_f1", "test_macro_f1"),
        ("Selected Test Macro F1", "selected_test_macro_f1", "test_macro_f1"),
        ("Selected Test Macro Precision", "selected_test_macro_precision", "test_macro_precision"),
        ("Selected Test Macro Recall", "selected_test_macro_recall", "test_macro_recall"),
        ("Selected Test Subset Accuracy", "selected_test_subset_accuracy", "test_subset_accuracy"),
    ]

    rows = []
    for display_name, row_key, metric_name in metric_specs:
        if not all(row_key in row for row in summary_rows):
            continue
        best_row = (
            min(summary_rows, key=lambda row: row[row_key])
            if is_lower_better(metric_name)
            else max(summary_rows, key=lambda row: row[row_key])
        )
        rows.append(
            [
                display_name,
                "Lower" if is_lower_better(metric_name) else "Higher",
                *[format_metric_value(metric_name, row[row_key]) for row in summary_rows],
                best_row["label"],
            ]
        )
    return headers, rows


def choose_compact_curve_metrics(metrics):
    chosen = []
    for metric_name in COMPACT_CURVE_PRIORITY:
        if metric_name in metrics and metric_name not in chosen:
            chosen.append(metric_name)
    for metric_name in metrics:
        if metric_name not in chosen:
            chosen.append(metric_name)
        if len(chosen) >= 2:
            break
    return chosen[:2]


def choose_compact_per_class_payloads(payloads):
    payload_by_metric = {payload["metric"]: payload for payload in payloads}
    chosen = []
    for metric_name in COMPACT_PER_CLASS_PRIORITY:
        payload = payload_by_metric.get(metric_name)
        if payload:
            chosen.append(payload)
        if len(chosen) >= 2:
            break
    return chosen


def export_ppt(report_dir: Path, context: ComparisonContext):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPT export. Install it or use --skip-ppt."
        ) from exc

    def rgb(color):
        return RGBColor(*color)

    report_name = context.report_name
    title = context.title
    runs = context.runs
    metrics = context.metrics
    selected_metric_keys = context.selected_metric_keys
    macro_metric_keys = context.macro_metric_keys
    summary_rows = context.summary_rows
    varying_rows = context.varying_rows
    all_varying_keys = context.all_varying_keys
    metric_figure_paths = context.metric_figure_paths
    macro_figure_payload = context.macro_figure_payload
    per_class_figure_payloads = context.per_class_figure_payloads
    headline_insights = context.headline_insights
    conclusion_lines = context.conclusion_lines
    scenario = context.scenario
    compact_mode = context.compact_mode

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    def apply_slide_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color)

    def add_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color=COLOR_TEXT,
        bold=False,
        align=PP_ALIGN.LEFT,
        font_name=FONT_FAMILY,
    ):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        paragraph.space_after = 0
        paragraph.space_before = 0
        paragraph.line_spacing = 1.1
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return textbox

    def add_bullet_box(slide, left, top, width, height, bullets, font_size=16, color=COLOR_TEXT):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.12
            for run in paragraph.runs:
                run.font.name = FONT_FAMILY
                run.font.size = Pt(font_size)
                run.font.color.rgb = rgb(color)
        return textbox

    def add_card(slide, left, top, width, height, title_text, body_lines, accent=COLOR_BLUE):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLOR_CARD)
        shape.line.color.rgb = rgb(COLOR_BORDER)
        shape.line.width = Pt(1.0)
        accent_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(0.15),
            Inches(height),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rgb(accent)
        accent_bar.line.fill.background()
        add_textbox(slide, left + 0.28, top + 0.18, width - 0.42, 0.36, title_text, 18, bold=True)
        add_bullet_box(slide, left + 0.28, top + 0.62, width - 0.42, height - 0.78, body_lines, font_size=11)

    def add_metric_card(slide, left, top, width, height, label_text, value_text, accent=COLOR_BLUE):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLOR_CARD)
        shape.line.color.rgb = rgb(COLOR_BORDER)
        shape.line.width = Pt(1.0)
        add_textbox(slide, left + 0.18, top + 0.18, width - 0.36, 0.28, label_text, 12, color=COLOR_MUTED, bold=True)
        add_textbox(slide, left + 0.18, top + 0.55, width - 0.36, 0.48, value_text, 20, color=accent, bold=True)

    def add_slide_header(slide, title_text, subtitle_text="", section_label="Experiment Comparison"):
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.55),
            Inches(0.32),
            Inches(1.9),
            Inches(0.34),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb(COLOR_HEADER_FILL)
        pill.line.fill.background()
        add_textbox(slide, 0.72, 0.38, 1.6, 0.2, section_label, 10, color=COLOR_BLUE, bold=True)
        add_textbox(slide, 0.55, 0.72, 8.7, 0.45, title_text, 25, bold=True)
        if subtitle_text:
            add_textbox(slide, 0.56, 1.14, 11.9, 0.28, subtitle_text, 11, color=COLOR_MUTED)

        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.55),
            Inches(1.42),
            Inches(12.2),
            Inches(0.02),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb(COLOR_BORDER)
        rule.line.fill.background()

    def add_footer(slide, left_text, right_text=""):
        add_textbox(slide, 0.55, 7.05, 6.6, 0.2, left_text, 9, color=COLOR_MUTED)
        if right_text:
            add_textbox(slide, 9.35, 7.05, 3.4, 0.2, right_text, 9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)

    def style_table(table, font_size):
        for row_index, row in enumerate(table.rows):
            for cell in row.cells:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(COLOR_CARD if row_index % 2 == 1 else COLOR_BG)
                if row_index == 0:
                    cell.fill.fore_color.rgb = rgb(COLOR_HEADER_FILL)
                cell.text_frame.word_wrap = True
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.line_spacing = 1.05
                    for run in paragraph.runs:
                        run.font.name = FONT_FAMILY
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = rgb(COLOR_TEXT)
                        run.font.bold = row_index == 0

        for cell in table.rows[0].cells:
            cell.fill.fore_color.rgb = rgb(COLOR_HEADER_FILL)

    def set_table_column_widths(table, column_count, sticky_columns):
        total_width = 12.2
        sticky_width = 0.0
        if sticky_columns >= 1:
            first_width = min(3.0, max(2.1, total_width * 0.24))
            table.columns[0].width = Inches(first_width)
            sticky_width += first_width

        for column_index in range(1, sticky_columns):
            width = min(1.8, max(1.1, total_width * 0.12))
            table.columns[column_index].width = Inches(width)
            sticky_width += width

        remaining_columns = column_count - sticky_columns
        if remaining_columns <= 0:
            return

        remaining_width = max(2.8, total_width - sticky_width)
        per_width = remaining_width / remaining_columns
        for column_index in range(sticky_columns, column_count):
            table.columns[column_index].width = Inches(per_width)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_NAVY)

        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.0),
            Inches(0.0),
            Inches(13.333),
            Inches(0.24),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(COLOR_CYAN)
        accent.line.fill.background()

        title_font_size = 28
        title_height = 0.85
        if len(title) > 58:
            title_font_size = 24
            title_height = 1.15
        if len(title) > 78:
            title_font_size = 21
            title_height = 1.35
        add_textbox(slide, 0.72, 0.9, 12.0, title_height, title, title_font_size, color=(255, 255, 255), bold=True)
        add_textbox(
            slide,
            0.74,
            2.0,
            10.6,
            0.4,
            "Weekly experiment comparison deck generated from existing run artifacts",
            14,
            color=(226, 232, 240),
        )

        run_chip_top = 2.7
        for index, run in enumerate(runs[:4]):
            chip = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.74 + index * 3.03),
                Inches(run_chip_top),
                Inches(2.74),
                Inches(0.5),
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = rgb((30, 41, 59))
            chip.line.color.rgb = rgb((71, 85, 105))
            add_textbox(
                slide,
                0.94 + index * 3.03,
                run_chip_top + 0.12,
                2.36,
                0.2,
                run.label,
                12,
                color=(255, 255, 255),
                bold=True,
                align=PP_ALIGN.CENTER,
            )

        meta_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.72),
            Inches(4.0),
            Inches(6.1),
            Inches(2.25),
        )
        meta_box.fill.solid()
        meta_box.fill.fore_color.rgb = rgb((30, 41, 59))
        meta_box.line.color.rgb = rgb((71, 85, 105))
        add_textbox(slide, 1.0, 4.3, 5.5, 0.3, "Report scope", 16, color=(255, 255, 255), bold=True)
        add_bullet_box(
            slide,
            1.0,
            4.72,
            5.25,
            1.3,
            [
                f"Comparison scenario: {format_comparison_scenario(scenario)}",
                f"Reference run: {runs[0].label}",
                f"Compared runs: {len(runs)}",
                f"Shared training metrics: {', '.join(metric_display_name(metric) for metric in metrics[:3])}",
            ],
            font_size=12,
            color=(226, 232, 240),
        )
        add_footer(
            slide,
            f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            f"Report name: {report_name}",
        )

    def add_run_info_slides():
        run_cards = [(run.label, build_run_profile_lines(run, all_varying_keys)) for run in runs]
        shared_dataset_card = build_shared_dataset_card_content(runs)
        for group_index, run_group in enumerate(chunked(run_cards, 4), start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_background(slide, COLOR_BG)
            add_slide_header(
                slide,
                "Experiment Setup and Run Context",
                "Keep the first --run as the reference baseline for meeting comparisons.",
                section_label="Run Setup",
            )

            if len(run_group) <= 2:
                positions = [(0.55, 1.72), (6.85, 1.72)]
                card_width = 5.9
                card_height = 3.15
            elif len(run_group) == 3:
                positions = [(0.55, 1.72), (6.85, 1.72), (0.55, 4.1)]
                card_width = 5.9
                card_height = 2.52
            else:
                positions = [
                    (0.55, 1.7),
                    (6.85, 1.7),
                    (0.55, 4.25),
                    (6.85, 4.25),
                ]
                card_width = 5.9
                card_height = 2.18

            for (card_left, card_top), (card_title, body_lines) in zip(positions, run_group):
                add_card(slide, card_left, card_top, card_width, card_height, card_title, body_lines)

            if shared_dataset_card and group_index == 1 and len(run_group) == 3:
                dataset_title, dataset_lines = shared_dataset_card
                add_card(
                    slide,
                    6.85,
                    4.1,
                    5.9,
                    2.52,
                    dataset_title,
                    dataset_lines,
                    accent=COLOR_CYAN,
                )

            add_footer(
                slide,
                "Run cards surface the main hyperparameters and monitoring settings without crowding a table.",
                f"Page group {group_index}",
            )

    def add_highlight_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            "Results Overview",
            "Headline findings are phrased for direct use in a weekly update.",
            section_label="Overview",
        )
        add_bullet_box(slide, 0.75, 1.8, 6.35, 3.3, headline_insights, font_size=15)

        primary_metric = metrics[0]
        if is_lower_better(primary_metric):
            best_row = min(summary_rows, key=lambda row: row[f"final_{primary_metric}"])
        else:
            best_row = max(summary_rows, key=lambda row: row[f"final_{primary_metric}"])

        add_metric_card(
            slide,
            7.45,
            1.9,
            2.45,
            1.3,
            "Best primary metric",
            best_row["label"],
            accent=COLOR_GOOD,
        )
        add_metric_card(
            slide,
            10.1,
            1.9,
            2.45,
            1.3,
            metric_display_name(primary_metric),
            format_metric_value(primary_metric, best_row[f"final_{primary_metric}"]),
            accent=COLOR_BLUE,
        )
        add_metric_card(
            slide,
            7.45,
            3.45,
            2.45,
            1.3,
            "Runs compared",
            str(len(runs)),
            accent=COLOR_CYAN,
        )
        add_metric_card(
            slide,
            10.1,
            3.45,
            2.45,
            1.3,
            "Metrics tracked",
            str(len(metrics) + len(selected_metric_keys)),
            accent=COLOR_CYAN,
        )

        add_footer(slide, "This slide replaces generic 'key takeaways' text with meeting-ready comparison statements.")

    def add_model_comparison_slide():
        headers, rows = build_model_comparison_rows(runs)
        add_table_slides(
            "Model Comparison",
            "This page is keyed off model family, variant, position encoding, and initialization rather than off raw run names.",
            headers,
            rows,
            sticky_columns=2,
            max_total_columns=7,
            max_body_rows=7,
            font_size=11,
            footer_text="Use this page to frame baseline vs RoPE, ViT vs CNN, or scratch vs pretrained before discussing metrics.",
        )

    def add_macro_metrics_slide():
        if not macro_metric_keys or not macro_figure_payload:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            "Macro Metrics",
            "Macro F1, precision, and recall are grouped here for quick meeting discussion.",
            section_label="Macro Metrics",
        )
        slide.shapes.add_picture(str(macro_figure_payload["path"]), Inches(0.62), Inches(1.82), width=Inches(7.25))
        add_card(
            slide,
            8.25,
            1.98,
            4.0,
            2.9,
            "Macro readout",
            macro_figure_payload["summary_lines"] or ["Macro metrics are available but no summary line was generated."],
            accent=COLOR_CYAN,
        )
        add_footer(slide, "Macro metrics help when accuracy alone hides uneven class performance.")

    def add_compact_metrics_slide():
        headers, rows = build_compact_metric_rows(summary_rows)
        if not rows:
            return
        add_table_slides(
            "Key Metrics",
            "Compact weekly snapshot of the final accuracy and selected-checkpoint macro metrics.",
            headers,
            rows,
            sticky_columns=2,
            max_total_columns=6,
            max_body_rows=8,
            font_size=11,
            footer_text="Compact mode keeps only the highest-signal summary rows for projection-friendly reading.",
        )

    def add_compact_curve_slide():
        chosen_metrics = choose_compact_curve_metrics(metrics)
        if not chosen_metrics:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            "Training Curves",
            "Compact mode keeps the most useful convergence views for weekly discussion.",
            section_label="Curves",
        )

        slot_positions = [(0.55, 1.72), (6.55, 1.72)]
        for (slot_left, slot_top), metric_name in zip(slot_positions, chosen_metrics):
            best_row = (
                min(summary_rows, key=lambda row: row[f"final_{metric_name}"])
                if is_lower_better(metric_name)
                else max(summary_rows, key=lambda row: row[f"final_{metric_name}"])
            )
            slide.shapes.add_picture(
                str(metric_figure_paths[metric_name]),
                Inches(slot_left),
                Inches(slot_top),
                width=Inches(5.65),
            )
            add_textbox(slide, slot_left, 5.65, 5.65, 0.22, metric_display_name(metric_name), 13, bold=True)
            add_textbox(
                slide,
                slot_left,
                5.95,
                5.65,
                0.28,
                f"Best final value: {best_row['label']} ({format_metric_value(metric_name, best_row[f'final_{metric_name}'])})",
                10,
                color=COLOR_MUTED,
            )

        add_footer(slide, "In compact mode, trend slides are limited so the deck stays short.")

    def add_compact_per_class_slide():
        selected_payloads = choose_compact_per_class_payloads(per_class_figure_payloads)
        if not selected_payloads:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            "Per-Class Analysis",
            "Only the most interpretable classwise views are kept in compact mode.",
            section_label="Per-Class",
        )

        if len(selected_payloads) == 1:
            payload = selected_payloads[0]
            slide.shapes.add_picture(str(payload["path"]), Inches(0.72), Inches(1.82), width=Inches(7.4))
            add_card(
                slide,
                8.55,
                2.0,
                3.7,
                3.1,
                metric_display_name(payload["metric"]),
                payload["summary"],
                accent=COLOR_CYAN,
            )
        else:
            positions = [(0.55, 1.85), (6.6, 1.85)]
            for (slot_left, slot_top), payload in zip(positions, selected_payloads):
                slide.shapes.add_picture(str(payload["path"]), Inches(slot_left), Inches(slot_top), width=Inches(5.8))
                add_textbox(slide, slot_left, 5.2, 5.8, 0.22, metric_display_name(payload["metric"]), 12, bold=True)
                add_bullet_box(slide, slot_left, 5.48, 5.8, 1.0, payload["summary"][:2], font_size=10)

        add_footer(slide, "Per-class bars help explain which composition elements benefit from each positional choice.")

    def add_table_slides(title_text, subtitle_text, headers, rows, sticky_columns=1, max_total_columns=6, max_body_rows=8, font_size=None, footer_text=""):
        if len(headers) <= max_total_columns:
            column_groups = [list(range(len(headers)))]
        else:
            scrollable_indexes = list(range(sticky_columns, len(headers)))
            scrollable_chunk_size = max(1, max_total_columns - sticky_columns)
            column_groups = []
            for chunk in chunked(scrollable_indexes, scrollable_chunk_size):
                column_groups.append(list(range(sticky_columns)) + chunk)

        row_groups = list(chunked(rows, max_body_rows)) or [[]]
        total_slides = len(column_groups) * len(row_groups)
        slide_number = 0

        for row_group in row_groups:
            for column_group in column_groups:
                slide_number += 1
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                apply_slide_background(slide, COLOR_BG)
                numbered_title = title_text if total_slides == 1 else f"{title_text} ({slide_number}/{total_slides})"
                add_slide_header(slide, numbered_title, subtitle_text, section_label="Comparison Table")

                subset_headers = [headers[index] for index in column_group]
                subset_rows = [
                    [shorten(str(row[index]), width=40, placeholder="...") for index in column_group]
                    for row in row_group
                ]
                local_font_size = font_size or estimate_table_font_size(
                    column_count=len(subset_headers),
                    row_count=len(subset_rows),
                )

                shape = slide.shapes.add_table(
                    len(subset_rows) + 1,
                    len(subset_headers),
                    Inches(0.55),
                    Inches(1.72),
                    Inches(12.2),
                    Inches(4.95),
                )
                table = shape.table

                for column_index, header in enumerate(subset_headers):
                    table.cell(0, column_index).text = header

                for row_index, row_values in enumerate(subset_rows, start=1):
                    for column_index, value in enumerate(row_values):
                        table.cell(row_index, column_index).text = value

                effective_sticky_columns = min(sticky_columns, len(subset_headers))
                set_table_column_widths(table, len(subset_headers), effective_sticky_columns)
                row_height = Inches(4.95 / max(2, len(subset_rows) + 1))
                for row in table.rows:
                    row.height = row_height
                style_table(table, local_font_size)

                add_footer(slide, footer_text)

    def add_curve_slides():
        figure_payloads = []
        for metric in metrics:
            best_row = min(summary_rows, key=lambda row: row[f"final_{metric}"]) if is_lower_better(metric) else max(summary_rows, key=lambda row: row[f"final_{metric}"])
            figure_payloads.append(
                {
                    "title": metric_display_name(metric),
                    "path": metric_figure_paths[metric],
                    "summary": f"Best final value: {best_row['label']} ({format_metric_value(metric, best_row[f'final_{metric}'])})",
                }
            )

        for page_index, group in enumerate(chunked(figure_payloads, 2), start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_background(slide, COLOR_BG)
            add_slide_header(
                slide,
                f"Training Curves ({page_index}/{len(list(chunked(figure_payloads, 2)))})",
                "Curves use a consistent title pattern and highlight the final operating point.",
                section_label="Curves",
            )

            slot_positions = [(0.55, 1.72), (6.55, 1.72)]
            for (slot_left, slot_top), payload in zip(slot_positions, group):
                slide.shapes.add_picture(str(payload["path"]), Inches(slot_left), Inches(slot_top), width=Inches(5.65))
                add_textbox(slide, slot_left, 5.65, 5.65, 0.22, payload["title"], 13, bold=True)
                add_textbox(slide, slot_left, 5.95, 5.65, 0.28, payload["summary"], 10, color=COLOR_MUTED)

            add_footer(slide, "Use these pages for trend discussion; use tables for exact values.")

    def add_confusion_matrix_slides():
        confusion_runs = [run for run in runs if run.confusion_matrix_csv]
        for run in confusion_runs:
            figure_path = ensure_confusion_matrix_figure(report_dir / "figures", run)
            if not figure_path:
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_background(slide, COLOR_BG)
            add_slide_header(
                slide,
                f"Confusion Matrix | {run.model_display_name}",
                "Selected-checkpoint error analysis pulled from existing evaluation outputs.",
                section_label="Error Analysis",
            )
            slide.shapes.add_picture(str(figure_path), Inches(0.68), Inches(1.8), width=Inches(6.6))

            bullets = []
            if run.selected_epoch is not None:
                bullets.append(f"Selected checkpoint epoch: {run.selected_epoch}")
            if "test_acc" in run.selected_metrics:
                bullets.append(
                    f"Selected test accuracy: {format_metric_value('test_acc', run.selected_metrics['test_acc'])}"
                )
            if "test_macro_f1" in run.selected_metrics:
                bullets.append(
                    f"Selected test macro F1: {format_metric_value('test_macro_f1', run.selected_metrics['test_macro_f1'])}"
                )
            confusion_pairs = top_confusion_pairs(run.confusion_matrix_csv)
            if confusion_pairs:
                bullets.extend(
                    [
                        f"Top confusion: true {true_label} predicted as {pred_label} ({count} samples)."
                        for count, true_label, pred_label in confusion_pairs[:2]
                    ]
                )
            add_card(slide, 7.7, 1.95, 4.55, 3.75, "Reading notes", bullets or ["Confusion matrix figure available, but no summary bullets were derived."])
            add_footer(slide, "This page gives a direct bridge from aggregate metrics to class-level failure modes.")

    def add_per_class_slides():
        for payload in per_class_figure_payloads:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_background(slide, COLOR_BG)
            add_slide_header(
                slide,
                f"Per-Class Analysis | {metric_display_name(payload['metric'])}",
                "Classwise selected-checkpoint metrics are grouped across runs for easier error analysis.",
                section_label="Per-Class",
            )
            slide.shapes.add_picture(str(payload["path"]), Inches(0.62), Inches(1.85), width=Inches(7.55))
            add_card(slide, 8.45, 1.98, 3.8, 3.45, "Key readout", payload["summary"], accent=COLOR_CYAN)
            add_footer(slide, "Use this page to explain which classes move together and which remain difficult.")

    def add_conclusion_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            "Conclusion",
            "Short, reusable wrap-up slide for weekly meetings.",
            section_label="Close",
        )
        add_bullet_box(slide, 0.8, 1.95, 7.3, 3.2, conclusion_lines, font_size=16)
        add_card(
            slide,
            8.55,
            2.1,
            3.7,
            2.3,
            "Recommended next use",
            [
                "Keep the baseline run first in the command line.",
                "Append new runs with readable labels.",
                "Reuse this deck directly for group meetings.",
            ],
            accent=COLOR_GOOD,
        )
        add_footer(slide, "The goal is to export a presentable deck, not a raw artifact dump.")

    add_title_slide()
    add_run_info_slides()
    add_highlight_slide()

    if compact_mode:
        add_compact_metrics_slide()
        add_macro_metrics_slide()
        add_compact_curve_slide()
        add_compact_per_class_slide()
        add_conclusion_slide()
    else:
        add_model_comparison_slide()

        overview_headers, overview_rows = build_results_overview_rows(summary_rows, metrics)
        add_table_slides(
            "Result Snapshot",
            "Compact overview of the most important final metrics and selected checkpoint epoch.",
            overview_headers,
            overview_rows,
            sticky_columns=2,
            max_total_columns=8,
            max_body_rows=7,
            footer_text="This table is intentionally compact so it remains readable on projected slides.",
        )

        metric_headers, metric_rows = build_main_metric_rows(summary_rows, metrics, selected_metric_keys)
        add_table_slides(
            "Main Metrics Comparison",
            "Each row is a meeting-friendly metric statement rather than a raw logger dump.",
            metric_headers,
            metric_rows,
            sticky_columns=2,
            max_total_columns=6,
            max_body_rows=8,
            footer_text="Direction indicates whether higher or lower values are preferable.",
        )

        if macro_metric_keys:
            macro_headers, macro_rows = build_macro_metric_rows(summary_rows, macro_metric_keys)
            add_table_slides(
                "Macro Metrics Table",
                "Exact selected-checkpoint macro metrics for precision, recall, and F1.",
                macro_headers,
                macro_rows,
                sticky_columns=2,
                max_total_columns=6,
                max_body_rows=8,
                footer_text="This table complements the macro bar chart on the next page.",
            )
            add_macro_metrics_slide()

        if varying_rows:
            config_headers = ["Config"] + [run.label for run in runs]
            config_table_rows = [[key, *values] for key, values in varying_rows]
            add_table_slides(
                "Config Differences",
                f"Showing the top {len(varying_rows)} differing settings for experiment traceability.",
                config_headers,
                config_table_rows,
                sticky_columns=1,
                max_total_columns=5,
                max_body_rows=7,
                font_size=10,
                footer_text="Long config lists are paginated automatically to avoid overflow and tiny fonts.",
            )

        add_curve_slides()
        add_confusion_matrix_slides()
        add_per_class_slides()
        add_conclusion_slide()

    ppt_path = report_dir / f"{report_name}.pptx"
    prs.save(ppt_path)
    return ppt_path


def export_seed_summary_ppt(
    report_dir: Path,
    report_name: str,
    title: str,
    scenario: str,
    seed_summary: dict,
):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPT export. Install it or use --skip-ppt."
        ) from exc

    def rgb(color):
        return RGBColor(*color)

    def format_percent(value):
        return f"{float(value) * 100.0:.2f}%"

    summary_rows = seed_summary["summary_rows"]
    delta_rows = seed_summary["delta_rows"]
    metric_figures = seed_summary["metric_figures"]
    insights = seed_summary["headline_insights"]
    conclusions = seed_summary["conclusion_lines"]
    models = seed_summary["models"]
    seeds = seed_summary["seeds"]
    reference_model = seed_summary["reference_model"]
    metrics = seed_summary["metrics"]
    supplementary_per_class_slides = seed_summary.get("supplementary_per_class_slides", [])

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    def apply_slide_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color)

    def add_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color=COLOR_TEXT,
        bold=False,
        align=PP_ALIGN.LEFT,
        font_name=FONT_FAMILY,
    ):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        paragraph.space_after = 0
        paragraph.space_before = 0
        paragraph.line_spacing = 1.1
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return textbox

    def add_bullet_box(slide, left, top, width, height, bullets, font_size=16, color=COLOR_TEXT):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.12
            for run in paragraph.runs:
                run.font.name = FONT_FAMILY
                run.font.size = Pt(font_size)
                run.font.color.rgb = rgb(color)
        return textbox

    def add_card(slide, left, top, width, height, title_text, body_lines, accent=COLOR_BLUE):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLOR_CARD)
        shape.line.color.rgb = rgb(COLOR_BORDER)
        shape.line.width = Pt(1.0)
        accent_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(0.15),
            Inches(height),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rgb(accent)
        accent_bar.line.fill.background()
        add_textbox(slide, left + 0.28, top + 0.18, width - 0.42, 0.36, title_text, 18, bold=True)
        add_bullet_box(slide, left + 0.28, top + 0.62, width - 0.42, height - 0.78, body_lines, font_size=11)

    def add_slide_header(slide, title_text, subtitle_text="", section_label="Seed Summary"):
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.55),
            Inches(0.32),
            Inches(1.8),
            Inches(0.34),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb(COLOR_HEADER_FILL)
        pill.line.fill.background()
        add_textbox(slide, 0.72, 0.38, 1.45, 0.2, section_label, 10, color=COLOR_BLUE, bold=True)
        add_textbox(slide, 0.55, 0.72, 9.2, 0.45, title_text, 24, bold=True)
        if subtitle_text:
            add_textbox(slide, 0.56, 1.14, 11.9, 0.28, subtitle_text, 11, color=COLOR_MUTED)
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.55),
            Inches(1.42),
            Inches(12.2),
            Inches(0.02),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb(COLOR_BORDER)
        rule.line.fill.background()

    def add_footer(slide, left_text, right_text=""):
        add_textbox(slide, 0.55, 7.05, 7.6, 0.2, left_text, 9, color=COLOR_MUTED)
        if right_text:
            add_textbox(slide, 9.0, 7.05, 3.75, 0.2, right_text, 9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)

    def style_table(table, font_size):
        for row_index, row in enumerate(table.rows):
            for cell in row.cells:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(COLOR_CARD if row_index % 2 == 1 else COLOR_BG)
                if row_index == 0:
                    cell.fill.fore_color.rgb = rgb(COLOR_HEADER_FILL)
                cell.text_frame.word_wrap = True
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = FONT_FAMILY
                        run.font.size = Pt(font_size)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = rgb(COLOR_TEXT)

    def set_table_widths(table, widths):
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, COLOR_NAVY)
    add_textbox(slide, 0.72, 0.88, 11.6, 0.76, title, 28, color=(255, 255, 255), bold=True)
    add_textbox(
        slide,
        0.74,
        1.82,
        10.9,
        0.35,
        "Aggregate multi-seed comparison deck generated from summary artifacts.",
        14,
        color=(226, 232, 240),
    )
    add_card(
        slide,
        0.72,
        3.95,
        6.3,
        2.25,
        "Report scope",
        [
            f"Scenario: {format_comparison_scenario(scenario)}",
            f"Seeds: {', '.join(str(seed) for seed in seeds)}",
            f"Models: {', '.join(MODEL_DISPLAY_NAMES.get(model, model) for model in models)}",
            f"Reference: {MODEL_DISPLAY_NAMES.get(reference_model, reference_model)}",
        ],
        accent=COLOR_CYAN,
    )
    add_footer(slide, f"Generated from {seed_summary['source_report_name']}", f"Report name: {report_name}")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, COLOR_BG)
    add_slide_header(
        slide,
        "Seed-Summary Highlights",
        "Use this page to lead with the average-performance story before opening per-metric evidence.",
        section_label="Overview",
    )
    add_bullet_box(slide, 0.78, 1.85, 6.65, 3.55, insights, font_size=16)
    best_model = max(summary_rows, key=lambda item: item["test_acc_mean"])
    add_card(
        slide,
        7.75,
        1.95,
        4.1,
        2.45,
        "Best aggregate model",
        [
            best_model["model_label"],
            f"Mean test acc: {format_percent(best_model['test_acc_mean'])}",
            f"Std: {format_percent(best_model['test_acc_std'])}",
        ],
        accent=COLOR_GOOD,
    )
    add_footer(slide, "This page frames the meeting around cross-seed reliability instead of a single run.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, COLOR_BG)
    add_slide_header(
        slide,
        "Aggregate Metrics Table",
        "Mean/std, min/max, and selected-epoch summary across seeds.",
        section_label="Aggregate Table",
    )
    headers = ["Model", "Seeds", "Best Val Acc", "Test Acc", "Macro F1", "Best Epoch"]
    rows = []
    for row in summary_rows:
        rows.append(
            [
                row["model_label"],
                str(row["num_seeds"]),
                f"{format_percent(row['best_val_acc_mean'])} +- {format_percent(row['best_val_acc_std'])}",
                f"{format_percent(row['test_acc_mean'])} +- {format_percent(row['test_acc_std'])}",
                f"{format_percent(row['macro_f1_mean'])} +- {format_percent(row['macro_f1_std'])}",
                f"{float(row['best_epoch_mean']):.2f} +- {float(row['best_epoch_std']):.2f}",
            ]
        )
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.55), Inches(1.72), Inches(12.2), Inches(4.95)).table
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    for r, row_values in enumerate(rows, start=1):
        for c, value in enumerate(row_values):
            table.cell(r, c).text = value
    set_table_widths(table, [2.2, 0.9, 2.2, 2.2, 2.2, 1.8])
    style_table(table, 11)
    add_footer(slide, "The aggregate table is the source-of-truth slide for mean/std discussion.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, COLOR_BG)
    add_slide_header(
        slide,
        "Delta Vs Reference",
        "Aggregate deltas are measured against the selected baseline reference model.",
        section_label="Delta Table",
    )
    headers = ["Model", "Best Val Acc Delta", "Test Acc Delta", "Macro F1 Delta"]
    rows = []
    for row in delta_rows:
        rows.append(
            [
                row["model_label"],
                format_delta("best_val_acc", row[f"best_val_acc_delta_vs_{reference_model}"]),
                format_delta("test_acc", row[f"test_acc_delta_vs_{reference_model}"]),
                format_delta("macro_f1", row[f"macro_f1_delta_vs_{reference_model}"]),
            ]
        )
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.8), Inches(1.95), Inches(11.6), Inches(3.4)).table
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    for r, row_values in enumerate(rows, start=1):
        for c, value in enumerate(row_values):
            table.cell(r, c).text = value
    set_table_widths(table, [2.6, 2.9, 2.9, 2.9])
    style_table(table, 12)
    add_footer(slide, "Positive delta means improvement over the chosen reference model.")

    for metric in metrics:
        mean_std_key = f"{metric}_mean_std"
        by_seed_key = f"{metric}_by_seed"
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            metric_display_name(metric),
            "Aggregate mean/std bars are paired with per-seed trajectories for the same metric.",
            section_label="Metric View",
        )
        if mean_std_key in metric_figures:
            slide.shapes.add_picture(str(metric_figures[mean_std_key]), Inches(0.55), Inches(1.82), width=Inches(5.85))
        if by_seed_key in metric_figures:
            slide.shapes.add_picture(str(metric_figures[by_seed_key]), Inches(6.85), Inches(1.82), width=Inches(5.55))
        add_footer(slide, f"{metric_display_name(metric)} is shown both as aggregate error bars and seed-by-seed lines.")

    for payload in supplementary_per_class_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_background(slide, COLOR_BG)
        add_slide_header(
            slide,
            payload["title"],
            payload["subtitle"],
            section_label="Per-Class",
        )
        slide.shapes.add_picture(str(payload["path"]), Inches(0.62), Inches(1.85), width=Inches(7.55))
        add_card(slide, 8.45, 1.98, 3.8, 3.45, "Key readout", payload["summary"], accent=COLOR_CYAN)
        add_footer(slide, "These slides complement the aggregate summary with class-level evidence.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, COLOR_BG)
    add_slide_header(
        slide,
        "Conclusion",
        "Short wrap-up slide for weekly meetings based on multi-seed evidence.",
        section_label="Close",
    )
    add_bullet_box(slide, 0.82, 1.95, 7.2, 3.5, conclusions, font_size=16)
    add_card(
        slide,
        8.5,
        2.05,
        3.7,
        2.25,
        "Recommended next use",
        [
            "Use this deck for the 3-seed mean/std story.",
            "Use per-seed decks only for debugging or variance discussion.",
            "Keep future model families in the same summary-report pipeline.",
        ],
        accent=COLOR_GOOD,
    )
    add_footer(slide, "This deck is rendered by generate_comparison_report.py from seed-summary artifacts.")

    ppt_path = report_dir / f"{report_name}.pptx"
    prs.save(ppt_path)
    return ppt_path


def build_seed_summary_report_data(args, project_root: Path):
    manifest_path = resolve_summary_manifest_path(
        args.results_dir,
        args.summary_report,
        args.summary_manifest,
        experiment_name=args.experiment_name,
    )
    source_manifest = load_json(manifest_path)
    per_class_report = load_per_class_report_data(args, project_root)
    source_report_name = str(source_manifest.get("report_name") or manifest_path.parent.name)
    report_name = args.report_name or source_report_name
    if report_name == source_report_name:
        report_dir = manifest_path.parent
    else:
        report_dir = build_report_artifact_dirs(
            results_dir=args.results_dir,
            report_name=report_name,
            experiment_name=args.experiment_name,
        )["report_dir"]
    report_dir.mkdir(parents=True, exist_ok=True)

    aggregate_csv_path = resolve_summary_artifact_path(project_root, source_manifest.get("aggregate_csv"))
    delta_csv_path = resolve_summary_artifact_path(project_root, source_manifest.get("delta_csv"))
    summary_rows = [
        {key: parse_summary_value(value) for key, value in row.items()}
        for row in load_csv_dict_rows(aggregate_csv_path)
    ]
    delta_rows = [
        {key: parse_summary_value(value) for key, value in row.items()}
        for row in load_csv_dict_rows(delta_csv_path)
    ]

    metric_figures = {}
    for key, value in (source_manifest.get("figures") or {}).items():
        metric_figures[key] = resolve_summary_artifact_path(project_root, value)

    metrics = list(source_manifest.get("metrics", []))
    models = list(source_manifest.get("models", []))
    scenario = detect_seed_summary_scenario(models)
    title = args.title or source_manifest.get("title") or f"Seed Summary: {report_name.replace('_', ' ').title()}"
    headline_insights = list(source_manifest.get("headline_insights", []))
    reference_model = str(source_manifest.get("reference_model", models[0] if models else "unknown"))
    conclusion_lines = build_seed_summary_conclusion_lines(summary_rows, reference_model, metrics)

    return {
        "manifest_path": manifest_path,
        "source_manifest": source_manifest,
        "source_report_name": source_report_name,
        "report_name": report_name,
        "report_dir": report_dir,
        "title": title,
        "scenario": scenario,
        "summary_rows": summary_rows,
        "delta_rows": delta_rows,
        "metric_figures": metric_figures,
        "headline_insights": headline_insights,
        "conclusion_lines": conclusion_lines,
        "metrics": metrics,
        "models": models,
        "seeds": list(source_manifest.get("seeds", [])),
        "reference_model": reference_model,
        "per_class_report_manifest": str(per_class_report["manifest_path"]) if per_class_report else None,
        "supplementary_per_class_slides": per_class_report["slides"] if per_class_report else [],
    }


def build_report_context(args, run_specs, project_root: Path):
    report_name = args.report_name or make_default_report_name(run_specs)
    runs = [
        load_run_artifacts(
            args.results_dir,
            project_root,
            run_name,
            label,
            experiment_name=args.experiment_name,
        )
        for run_name, label in run_specs
    ]
    runs = ensure_unique_run_labels(runs)
    metrics = determine_metrics(runs, explicit_metrics=args.metrics)
    if not metrics:
        raise ValueError("No shared numeric metrics were found across the selected runs.")

    scenario = detect_comparison_scenario(runs)
    title = args.title or build_default_title(runs, scenario)
    selected_metric_keys = determine_selected_metric_keys(runs, metrics)
    macro_metric_keys = determine_macro_metric_keys(selected_metric_keys, metrics)
    summary_rows = build_summary_rows(runs, metrics, selected_metric_keys)
    selected_varying_keys, all_varying_keys = choose_varying_config_keys(
        runs, args.max_config_rows
    )
    varying_rows = build_varying_rows(runs, selected_varying_keys, args.max_config_rows)

    return ComparisonContext(
        report_name=report_name,
        title=title,
        scenario=scenario,
        compact_mode=args.compact,
        runs=runs,
        metrics=metrics,
        selected_metric_keys=selected_metric_keys,
        macro_metric_keys=macro_metric_keys,
        summary_rows=summary_rows,
        varying_rows=varying_rows,
        all_varying_keys=all_varying_keys,
        metric_figure_paths={},
        macro_figure_payload=None,
        per_class_figure_payloads=[],
        headline_insights=[],
        conclusion_lines=[],
    )


def main():
    args = parse_args()
    project_root = Path.cwd()
    if args.runs:
        run_specs = [parse_run_spec(spec) for spec in args.runs]
        context = build_report_context(args, run_specs, project_root)
        report_dir, figures_dir = ensure_report_dirs(
            args.results_dir,
            context.report_name,
            experiment_name=args.experiment_name,
            runs=context.runs,
        )

        for metric in context.metrics:
            context.metric_figure_paths[metric] = plot_metric(figures_dir, context.runs, metric)

        if context.macro_metric_keys:
            context.macro_figure_payload = plot_macro_metrics(
                figures_dir, context.runs, context.macro_metric_keys
            )

        per_class_metric_keys = determine_available_per_class_metric_keys(context.runs)
        for metric_name in per_class_metric_keys:
            figure_path, class_names = plot_per_class_metric(figures_dir, context.runs, metric_name)
            context.per_class_figure_payloads.append(
                {
                    "metric": metric_name,
                    "path": figure_path,
                    "summary": summarize_per_class_metric(context.runs, metric_name, class_names),
                }
            )

        context.headline_insights = build_headline_insights(
            context.runs, context.summary_rows, context.metrics, context.scenario
        )
        context.conclusion_lines = build_conclusion_lines(
            context.runs,
            context.summary_rows,
            context.metrics,
            per_class_metric_keys,
            context.scenario,
        )

        summary_csv_path, config_csv_path, overview_path, presentation_summary_path, manifest_path = save_summary_outputs(
            report_dir=report_dir,
            title=context.title,
            scenario=context.scenario,
            runs=context.runs,
            metrics=context.metrics,
            selected_metric_keys=context.selected_metric_keys,
            summary_rows=context.summary_rows,
            varying_keys=context.all_varying_keys,
            headline_insights=context.headline_insights,
            conclusion_lines=context.conclusion_lines,
        )
        publication_paths = {
            "selected_checkpoint_table": write_publication_selected_checkpoint_table(
                report_dir, context.runs, context.summary_rows, context.selected_metric_keys
            ),
            "selected_test_accuracy": plot_selected_test_accuracy_summary(
                figures_dir, context.runs, context.summary_rows
            ),
        }
        captions_path = write_publication_captions(
            report_dir=report_dir,
            title=context.title,
            runs=context.runs,
            metrics=context.metrics,
            selected_metric_keys=context.selected_metric_keys,
            publication_paths=publication_paths,
        )

        print(f"Report directory: {report_dir}")
        print(f"Summary CSV: {summary_csv_path}")
        print(f"Config CSV: {config_csv_path}")
        print(f"Overview Markdown: {overview_path}")
        print(f"Presentation Summary JSON: {presentation_summary_path}")
        print(f"Manifest JSON: {manifest_path}")
        print(f"Publication selected-checkpoint table: {publication_paths['selected_checkpoint_table']}")
        print(f"Publication captions: {captions_path}")
        for metric, path in context.metric_figure_paths.items():
            print(f"Figure ({metric}): {path}")
        if publication_paths["selected_test_accuracy"]:
            print(f"Publication Figure (selected test accuracy): {publication_paths['selected_test_accuracy']}")
        if context.macro_figure_payload:
            print(f"Macro Figure: {context.macro_figure_payload['path']}")
        for payload in context.per_class_figure_payloads:
            print(f"Per-class Figure ({payload['metric']}): {payload['path']}")

        if not args.skip_ppt:
            ppt_path = export_ppt(report_dir=report_dir, context=context)
            print(f"PPTX: {ppt_path}")
        return

    seed_summary = build_seed_summary_report_data(args, project_root)
    presentation_summary_path, manifest_path = save_seed_summary_report_outputs(
        report_dir=seed_summary["report_dir"],
        title=seed_summary["title"],
        scenario=seed_summary["scenario"],
        source_manifest_path=seed_summary["manifest_path"],
        source_manifest=seed_summary["source_manifest"],
        headline_insights=seed_summary["headline_insights"],
        conclusion_lines=seed_summary["conclusion_lines"],
        per_class_manifest_path=Path(seed_summary["per_class_report_manifest"]) if seed_summary.get("per_class_report_manifest") else None,
    )

    print(f"Report directory: {seed_summary['report_dir']}")
    print(f"Source summary manifest: {seed_summary['manifest_path']}")
    print(f"Presentation Summary JSON: {presentation_summary_path}")
    print(f"Manifest JSON: {manifest_path}")
    for key, path in seed_summary["metric_figures"].items():
        print(f"Figure ({key}): {path}")

    if not args.skip_ppt:
        ppt_path = export_seed_summary_ppt(
            report_dir=seed_summary["report_dir"],
            report_name=seed_summary["report_name"],
            title=seed_summary["title"],
            scenario=seed_summary["scenario"],
            seed_summary=seed_summary,
        )
        print(f"PPTX: {ppt_path}")


if __name__ == "__main__":
    main()
